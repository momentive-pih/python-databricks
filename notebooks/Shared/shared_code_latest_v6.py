# Databricks notebook source
# -*- coding: utf-8 -*-
"""
Created on Mon Feb  3 16:48:39 2020

@author: 809917
"""
#**************************************************
#importing required packages
#**************************************************
import glob
import pytesseract as pyt
import logging
import datetime
import configparser
import shutil
import fitz
import os
from PIL import Image, ImageFilter
from wand.image import Image as wimage
import PIL
import docx
import pptx
import pandas as pd
import PyPDF2 
import openpyxl
import csv
import re
import nltk
from nltk import ngrams
from outlook_msg import Message
import json
import pyodbc
import camelot
import numpy as np
import shutil
import email
from datetime import datetime
from dateutil import parser as date_parser
config = configparser.ConfigParser()
#This configuration path should be configured in Blob storage
config.read("/dbfs/mnt/momentive-configuration/config-file.ini")


#Loging environment setup
logger = logging.getLogger('momentive')
logger.setLevel(logging.DEBUG)
fh = logging.FileHandler("shared_main_code2.log", 'w')
fh.setLevel(logging.DEBUG)
ch = logging.FileHandler("shared_main_code_error2.log", 'w')
ch.setLevel(logging.ERROR)
formatter =logging.Formatter(fmt = '%(asctime)s %(levelname)-8s %(message)s', datefmt='%Y-%m-%d %H:%M:%S')
fh.setFormatter(formatter)
ch.setFormatter(formatter)
logger.addHandler(ch)
logger.addHandler(fh)
text_folder_list = []
image_folder_list = []
record_folder_list = []
valid_folder_list =[]
folder_list =[]
extracted_file_list = []
native_scanned_folder_list =[]

#****************************************************
#function name: path_exists
#Objective: To empty or create the folders
#****************************************************
def path_exists(file_path):
  try:
    print('path_exists',file_path)
    logger.info("Executing path_exists function to create new folder in the {}".format(file_path))
    #dbutils.fs.rm(file_path.replace("/dbfs",""),True)
    if os.path.exists(file_path):
      #print('kamal')
      shutil.rmtree(file_path.replace('//','/'))
    #print(file_path)
    dbutils.fs.mkdirs(file_path.replace("/dbfs","dbfs:").replace('//','/'))
  except Exception as e:
    logger.error("Error in path_exists function while creation of : ".format(file_path),exc_info=True)

def unstructure_processed_data(unstructure_processed_data_query,category,product_type,product,data_extract,is_relevant,sql_conn,cursor):
    insert_query = \
    unstructure_processed_data_query.format(category,product_type,product,data_extract,'getdate()','getdate()',is_relevant)
    update_operation(insert_query,sql_conn,cursor) 

def excel_date(content,excel_file):
  flag=0
  date =None
  rex1 = re.compile(r'\d{1,2}\s*\/\d{1,2}\s*\/\d{4}') #12/12/2002
  rex2 = re.compile(r'[a-zA-Z]*\s*\d{1,2}\s*,\s*\d{4}') #Jan 23, 2002
  rex3 = re.compile(r'\d{1,2}\s*\-\s*[a-zA-Z]*\s*\-\s*\d{4}') #12-Jan-2002
  rex4 = re.compile(r'\d{1,2}\s*\-\s*[a-zA-Z]*\s*\-\s*\d{2}') #12-Jan-02
  rex5 = re.compile(r'\d{1,4}\s*\-\d{1,2}\s*\-\d{2}') #2002-12-12rex5 = re.compile(r'\d{1,4}\s*\-\d{1,2}\s*\-\d{2}') #2002-12-12
  rex6 = re.compile(r'\d{1,2}\s*\-\d{1,2}\s*\-\d{4}') #2002-12-12
  rex7 = re.compile(r'\d{1,4}\s*\-\d{1,2}\s*\-\d{2}')
  #for text in text_files:
  #content = open(text, 'r').read() 
  #month =['jan','feb','mar','apr','may','jun']
 # print(content)
  rex_text = rex4.findall(content) 
  rex_text1 = rex1.findall(content)
  rex_text2 = rex2.findall(content)
  rex_text3 = rex3.findall(content)
  rex_text4 = rex5.findall(content)
  rex_text5 = rex6.findall(content)
  rex_text6 = rex7.findall(content)
  if len(rex_text)>0:
    #rex_text= rex_text.replace(' ','')
    try:
      rex_text.sort(key = lambda date: datetime.strptime( date.replace(' ',''), '%d-%b-%y').date(), reverse=True) 
      date = rex_text[0]
     # print(date)
      flag=1
    except Exception as e:
      logger.error(e,exc_info=True)
  if len(rex_text1)>0:
   # rex_text1= rex_text1.replace(' ','')
    try:
      rex_text1.sort(key = lambda date: datetime.strptime(date.replace(' ',''), '%d/%m/%Y').date(), reverse=True) 
      date = rex_text1[0]
      #print(date)
      flag=1
    except Exception as e:
      logger.error(e,exc_info=True)
#   elif len(rex_text1)>0:
#     rex_text1.sort(key = lambda date: datetime.strptime(date, '%d/%b/%Y').date(), reverse=True) 
#     date = rex_text1[0]
#     print(date)
#     flag=1
  if len(rex_text2)>0:
    try:
      #rex_text2= rex_text2.replace(' ','')
      rex_text2.sort(key = lambda date: datetime.strptime(date.replace(' ',''), '%b %d,%Y').date(), reverse=True) 
      
      date = rex_text2[0]
      #print(date)
      flag=1
    except Exception as e:
      logger.error(e,exc_info=True)
  if len(rex_text3)>0:
    #rex_text3= rex_text3.replace(' ','')
    try:
      rex_text3.sort(key = lambda date: datetime.strptime(date.replace(' ',''), '%d-%b-%Y').date(), reverse=True) 
      date = rex_text3[0]
    #  print(date)
      flag=1
    except Exception as e:
      logger.error(e,exc_info=True)
  if len(rex_text4)>0:
    #rex_text4= rex_text4.replace(' ','')
    try:
      rex_text4.sort(key = lambda date: datetime.strptime(date.replace(' ',''), '%Y-%m-%d').date(), reverse=True) 
      date = rex_text4[0]
    #  print(date)
      flag=1
    except Exception as e:
      logger.error(e,exc_info=True)
  if len(rex_text5)>0:
    #rex_text5= rex_text5.replace(' ','')
    try:
      rex_text5.sort(key = lambda date: datetime.strptime(date.replace(' ',''), '%d-%m-%Y').date(), reverse=True) 
      date = rex_text5[0]
     # print(date)
      flag=1
    except Exception as e:
      logger.error(e,exc_info=True)
  if len(rex_text6)>0:
    #rex_text5= rex_text5.replace(' ','')
    try:
      rex_text6.sort(key = lambda date: datetime.strptime(date.replace(' ',''), '%Y-%m-%d').date(), reverse=True) 
      date = rex_text6[0]
    #  print(date)
      flag=1
    except Exception as e:
      logger.error(e,exc_info=True)
#   else:
#     flag=0

  if flag==0:
    try:
      head, tail = os.path.split(excel_file)
      file_name = tail.rsplit('.',1)[0]
      rex = rex3.findall(str(file_name)) or rex6.findall(str(file_name)) or rex1.findall(str(file_name)) or rex2.findall(str(file_name))\
            or rex4.findall(str(file_name)) or rex5.findall(str(file_name))
      if rex:
        date = str(rex[0])
    except Exception as e:
      logger.error(e,exc_info=True)
  return  date

def key_data_extract_external_source(valid_path,sql_conn,cursor,category,unstructure_processed_data_query,excel_date_found):
  global key_value_df_master_data
  
  try:
    logger.info('Executing key_data_extract_external_source function to load the data in sql server for {} excel'.format(valid_path))
    if os.path.exists(valid_path + 'relevant_data_files/'):
      files = glob.glob(valid_path + 'relevant_data_files/' + '*.csv')
      for file in files:
        json_list = []
     #   print(file)
        non_rel_data = pd.read_csv(file, encoding='iso-8859-1')
        product = 'Product'
        temp_data = non_rel_data.copy()
        temp_data.drop([product, 'Product_category', 'Component', 'is_relevant'], axis=1, inplace=True)
        df_dict = temp_data.to_json(orient='records', lines=False, force_ascii=False)
        d = json.loads(df_dict)
        print('d',len(d))
        for i in range(len(d)):
          if 'date' not in d[i] and 'Date' not in d[i]:
            d[i]['date'] = excel_date_found
          b = json.dumps(d[i], ensure_ascii=False)
          json_list.append(b)
        #key_value_df = pd.DataFrame(json_list, columns =['values']) 
        print(len(json_list))
        print(non_rel_data.shape)
        non_rel_data['values'] = json_list
        key_value_df_master = non_rel_data
        key_value_df_master_data = key_value_df_master.loc[:, ['Product_category', product, 'values', 'is_relevant']]
        #print()
        for i in range(key_value_df_master_data.shape[0]):
          if str(key_value_df_master_data['Product_category'][i]).strip().lower() == 'matnbr':                             
              product_value = '{0:0>18}'.format(str(key_value_df_master_data['Product'][i]).strip())
          else:
            product_value = str(key_value_df_master_data['Product'][i])
          if str(key_value_df_master_data['Product_category'][i]).strip().lower() == 'nan':                             
              product_cat_value = 'null'
          else:
            product_cat_value = str(key_value_df_master_data['Product_category'][i])
          unstructure_processed_data(unstructure_processed_data_query,category,product_cat_value
                                     ,product_value,key_value_df_master_data['values'][i].replace("'", 
                                     "''"),key_value_df_master_data['is_relevant'][i],sql_conn,cursor)
          
          
  except Exception as e:
    logger.error("Error in key_data_extract_external_source function while loading data in sql server  {} \
                 excel".format(valid_path),exc_info=True)
    
def data_validation_to_relevant_non_relevant_split(data_delta, valid_path, primary_column, comp,sql_conn,cursor,product_inscope_df):
  global consol_data,final
  try:
    #sql_conn = SQL_connection()
    logger.info('Executing data_validation_to_relevant_non_relevant_split function for {} excel'.format(valid_path))
#     query_product_list = 'select * from [momentive].[product_inscope]'
#     product_inscope = pd.read_sql(query_product_list, sql_conn)
#     CAS_list = list(product_inscope.cas_no)

#     nam_prod_list = list(product_inscope.nam_prod)
#     bdt_list = list(product_inscope.bdt)
    
    regex1 = re.compile(r'(\d+-\d+-\d+)', re.I) #CAS number formatting
    regex2 = re.compile(r'(\w+-\d{5})', re.I)  #Y-Number formatting
    regex3 = re.compile(r'(.+/+.+)', re.I) #
    
    reg_ex = [] 
    reg_ex1 = []
    reg_ex2=[]
    data_delta1 =[]
    dbutils.fs.rm((valid_path +'relevant_data_files/').replace("/dbfs",""),True)
  #  print('data_delta.shape',data_delta.shape)
    if not data_delta.shape[0]==0:
      for i in range(data_delta.shape[0]):         
        if not pd.isnull(data_delta.loc[i, primary_column]):
          product = data_delta.loc[i, primary_column]
       #   print('lll',str(product))
          reg_ex = regex1.findall(str(product)) or regex2.findall(str(product)) or regex3.findall(str(product))
          reg_ex1 = regex1.findall(str(product))# and regex2.findall(str(product))
          reg_ex2 = regex1.findall(str(product))
         # print('reg_ex',reg_ex)
          #print('reg_ex1',reg_ex1)
          if len(reg_ex1)>0:            
            for reg_len in  reg_ex1:
                data_delta.loc[i, primary_column]=reg_len.strip()
                data_delta1.append(data_delta.iloc[i].values.tolist())
          
          elif len(reg_ex)>0:           
            if '/' in reg_ex[0]:
                for reg_len in reg_ex[0].split('/') :
                  data_delta.loc[i, primary_column]=reg_len.strip()
                  data_delta1.append(data_delta.iloc[i].values.tolist())
            else:  
                data_delta.loc[i, primary_column]=reg_ex[0].strip()  
                data_delta1.append(data_delta.iloc[i].values.tolist())
          else:
            if len(reg_ex2)>0:  
              for reg_len in reg_ex2 :
                  data_delta.loc[i, primary_column]=reg_len.strip()
                  data_delta1.append(data_delta.iloc[i].values.tolist())
            else:
              data_delta1.append(data_delta.iloc[i].values.tolist())
        #data_delta =  data_delta.append(data_copy, ignore_index=True)
      
      data_delta1 = pd.DataFrame(data_delta1)
      data_delta1.columns =data_delta.columns
      data_delta = data_delta1 
   #   data_delta.to_csv(valid_path + 'cal-prop.csv')
   #   print(data_delta)
   #   print(data_delta[primary_column].values.tolist())
 
      master_relevant = data_delta.copy()
      master_relevant.rename(columns = {primary_column:'Product'}, inplace=True)
      cas_df = data_delta[primary_column].isin(product_inscope_df[product_inscope_df['Type'].str.contains('NUMCAS')]
                         ['Text'].values.tolist())                                
      
      cas_final = data_delta[cas_df]
      cas_final['Product_category'] = 'NUMCAS'
      nam_prod_list_df = data_delta[primary_column].isin(product_inscope_df[product_inscope_df['Type'].str.contains('NAMPROD')]
                         ['Text'].values.tolist())
      nam_prod_final = data_delta[nam_prod_list_df]
      nam_prod_final['Product_category'] = 'NAMPROD'
      bdt_df = data_delta[primary_column].isin(product_inscope_df[product_inscope_df['Type'].str.contains('BDT')]
                         ['Text'].values.tolist())
      bdt_final = data_delta[bdt_df]
      bdt_final['Product_category'] = 'BDT'
      material_no_copy=data_delta[primary_column].copy()
      material_no_copy.columns=primary_column
      #print('material_no_copy.columns',material_no_copy.columns)
      
      material_no_copy = material_no_copy.apply(lambda x: '{0:0>18}'.format(x))
     # print(material_no_copy)
      MATNBR_df = material_no_copy.isin(product_inscope_df[product_inscope_df['Type'].str.contains('MATNBR')]
                         ['Text'].values.tolist())
      
      MATNBR_final = data_delta[MATNBR_df]
      MATNBR_final['Product_category'] = 'MATNBR'
      consol_data = pd.concat([cas_final, nam_prod_final, bdt_final,MATNBR_final])
      consol_data.rename(columns = {primary_column:'Product'}, inplace=True)
      consol_data['Component'] = comp
      consol_data['is_relevant'] = 1
      master_consol_data = consol_data.copy()
      master_consol_data.drop(columns={'Product_category', 'Component', 'is_relevant'}, inplace=True)
      #print('valid_path',valid_path)
      #print('consol_data.shape',consol_data.shape)
      
      #dbutils.fs.mkdirs((valid_path +'relevant_data_files/').replace("/dbfs","dbfs:")) 
      path_exists(valid_path +'relevant_data_files/')
      if not consol_data.shape[0]==0:
        if os.path.exists(valid_path + 'relevant_data_files/'):
          print('oooo',valid_path + 'relevant_data_files/')
        #relevant_files = consol_data.to_csv(valid_path + 'relevant_data_files/' + 'relevant_data' +'.csv', index=None, header=True)
        consol_data.to_csv(valid_path + 'relevant_data_files/' + 'relevant_data.csv', index=None, header=True)
      final = master_relevant.append(master_consol_data)
      final.drop_duplicates(keep=False, inplace=True)
      final.reset_index(drop=True, inplace=True)
      final['Component'] = comp
      final['is_relevant'] = 0
      final['Product_category'] = np.nan
      #print('final.shape',final.shape)
      final['Product_category'].fillna("null", inplace = True) 
      if not final.shape[0]==0:            
        final.to_csv(valid_path + 'relevant_data_files/' + 'non_relevant_data.csv', index=None, header=True)   
  except Exception as e:
    logger.error("Error in data_validation_to_relevant_non_relevant_split function while loading relavent and non-relavent data for {} \
                 excel".format(valid_path),exc_info=True)
    
def excel_full_delta_load(valid_path, relevant_data,sql_conn,cursor,sheet_name):
  print('relevant_data', relevant_data.shape)
  global data_delta
  data_delta = pd.DataFrame()
  try:
    logger.info('Executing excel_full_delta_load function for {}'.format(valid_path))
    print('vliadjjjjjjjjj',valid_path + sheet_name+ '/')
    if not os.path.exists(valid_path + sheet_name+ '/'):
      print('lllllllooooooo',valid_path + 'valid/' + '*.csv')
      dbutils.fs.mkdirs((valid_path +sheet_name+ '/').replace("/dbfs","dbfs:")) 
      relevant_data.to_csv(valid_path +sheet_name+ '/' + 'valid_data.csv', index=None, header=True, encoding='iso-8859-1')
      match_data = pd.DataFrame()
      flag=1
    else:
      print('lllllll',valid_path + sheet_name+ '/')
      mat = glob.glob(valid_path + sheet_name+ '/' + '*.csv')
      for m in mat:
        match_data = pd.read_csv(m, encoding='iso-8859-1')
        flag=0
#        match_data.to_csv(valid_path + 'valid/' + 'valid_data1.csv')
    match_data.replace({r'[^\x00-\x7F]+':''}, regex=True, inplace=True)
#    print('relevant_data count', relevant_data.shape[0])
    #print()
    print(match_data.columns)
    print(relevant_data.columns)
    print('match_data_count', match_data.shape)
    #list(df1.align(df2)), ignore_index=True
    #data_delta = pd.concat(list(relevant_data.align(match_data)),ignore_index=True)
    data_delta = relevant_data.append(match_data)
    print('data_delta', data_delta.shape[0])
    data_delta.drop_duplicates(keep=False, inplace=True)
    data_delta.reset_index(drop=True, inplace=True)
   # data_delta1 = pd.concat(list(data_delta.align(match_data)),ignore_index=True)
    data_delta1 = data_delta.append(match_data)
    dup = data_delta1.duplicated(keep='first')
    data_delta2 = data_delta1[dup]
    data_delta2.reset_index(drop=True, inplace=True)
#    data_delta2.reset_index(drop=True, inplace=True)
 #   print('data_delta2', data_delta2.shape[0])

    if not data_delta.shape[0]==0 and not flag==1:
      data_to_valid = pd.read_csv(valid_path + sheet_name+ '/' + 'valid_data.csv', encoding='iso-8859-1')
      data_to_m = data_delta.append(data_delta2)
      data_to_m.drop_duplicates(keep=False, inplace=True)
      data_to_m.reset_index(drop=True, inplace=True)
      data_to_v = pd.concat(list(data_to_valid.align(data_to_m)),ignore_index=True) 
   #   data_to_v = data_to_valid.append(data_to_m)
      data_to_v.reset_index(drop=True, inplace=True)
      dbutils.fs.rm((valid_path +sheet_name+ '/').replace("/dbfs",""),True)
      data_to_v.to_csv(valid_path + sheet_name+ '/' + 'valid_data.csv', index=None, encoding='iso-8859-1')
      data_delta = data_to_m.copy()
#      print('data delta after deletion', data_delta.shape[0])
   
    if not data_delta2.shape[0]==0 and not flag==1:
      #data_to_v = pd.concat(list(data_to_valid.align(data_delta2)),ignore_index=True) 
      data_to_v = data_to_valid.append(data_delta2)
      data_to_v.drop_duplicates(keep=False, inplace=True)
      data_to_v.reset_index(drop=True, inplace=True)
      dbutils.fs.rm((valid_path +sheet_name+ '/').replace("/dbfs",""),True)
      data_to_v.to_csv(valid_path + sheet_name+ '/' + 'valid_data.csv', index=None, encoding='iso-8859-1')
    return data_delta, valid_path
  except Exception as e:
    logger.error("Error in excel_full_delta_load function while loading data from {} ".format(valid_path),exc_info=True)
    
def reading_excel_data_from_source(valid_path, files, component_data, primary_column, comp,sql_conn,cursor,product_inscope_df,unstructure_processed_data_query,excel_date_found,category,sheet_name):
  global relevant_data  
  try:
    logger.info('Executing reading_excel_data_from_source function for {} sheet in {}'.format(files,valid_path))
    component_columns = list(set(component_data['column_name']))
    component_columns1=[]
    for comp_col in component_columns:
      component_columns1.append(comp_col.strip())
    component_columns = list(set(component_columns1))
    #print('component_columns',component_columns)
    data_valid_extract = pd.read_csv(files, encoding='iso-8859-1', header=None)
    data_valid_extract = data_valid_extract.dropna(how='all',axis=0)
    data_valid_extract.reset_index(drop=True, inplace=True)
    #print('data_valid_extract.shape',data_valid_extract.shape)
    for i in range(data_valid_extract.shape[0]):
      row_list = list(data_valid_extract.loc[i,:])
      start_row_count = list(set(row_list) & set(component_columns))
      if len(start_row_count) >=3:
         value_of_column = i
    valid_data = data_valid_extract[int(value_of_column):]
    valid_data = valid_data.rename(columns=valid_data.iloc[0])
    valid_data.drop(valid_data.index[0], inplace=True)
    valid_data.reset_index(drop=True, inplace=True)
    valid_data.columns = valid_data.columns.str.replace('\n',' ')
    valid_data.columns = valid_data.columns.str.strip()
    valid_data.columns = valid_data.columns.str.replace(r'[^\x00-\x7F]+', '')
    relevant_data = valid_data.loc[:, component_columns]
    relevant_data.replace({r'[^\x00-\x7F]+':''}, regex=True, inplace=True)
    relevant_data.drop_duplicates(keep='first', inplace=True)
    relevant_data.reset_index(drop=True, inplace=True)
    relevant_data = relevant_data.loc[:,~relevant_data.columns.duplicated()]
    data_delta, valid_path = excel_full_delta_load(valid_path, relevant_data,sql_conn,cursor,sheet_name)
    #print('data_delta',data_delta)
   # print('valid_path',valid_path)
    data_validation_to_relevant_non_relevant_split(data_delta, valid_path, primary_column, comp,sql_conn,cursor,product_inscope_df)
    key_data_extract_external_source(valid_path,sql_conn,cursor,category,
                      unstructure_processed_data_query,excel_date_found)
  except Exception as e:
    logger.error("Error in reading_excel_data_from_source function while extracting from {} sheet in the {} server"  
                 .format(files,valid_path),exc_info=True) 

def reading_excel_sources(source_type, sql_conn,cursor):
  try:
    logger.info('Executing reading_excel_sources function for {}'.format(source_type))
    excel_momentive_source = config.get('mount_path','external_excel_source')
    data_excel_external_source = pd.read_sql(excel_momentive_source, sql_conn)
    dataframe_excel_sources = data_excel_external_source[(data_excel_external_source['source_type']==source_type) & \
                                                         (data_excel_external_source['is_active_folder'].astype('str')=='1') & \
                                                         (data_excel_external_source['is_active_column'].astype('str')=='1') & \
                                                         (data_excel_external_source['is_active_sheet'].astype('str')=='1')]
    #print(dataframe_excel_sources['sheet_name'])
    primary_field = data_excel_external_source[(data_excel_external_source['source_type']==source_type) & \
                                               (data_excel_external_source['is_active_folder'].astype('str')=='1') & \
                                               (data_excel_external_source['is_active_column'].astype('str')=='1') & \
                                               (data_excel_external_source['is_active_sheet'].astype('str')=='1') & \
                                               (data_excel_external_source['is_primary'].astype('str')=='1')]
    
    primary_col = list(set(primary_field.column_name.values))
    external_sheet = list(pd.unique(dataframe_excel_sources['sheet_name']))
    #print(external_sheet)
    return dataframe_excel_sources, external_sheet, primary_col
  except Exception as e:
    logger.error("Error in reading_excel_sources function while reading {} in sql server".format(source_type),exc_info=True)    
    
def excel_extract2_key_value_pair(valid_path, sql_conn,cursor,category,product_inscope_df,unstructure_processed_data_query,excel_date_found):
  try:
    logger.info('Executing excel_extract2_key_value_pair function for {}'.format(valid_path))
    external_source_data = config.get('mount_path','external_excel_source')
    source_type_valid = pd.read_sql(external_source_data, sql_conn)
    list_components = list(pd.unique(source_type_valid.source_type))
    #print('list_components',list_components)
    if not list_components:
      pass
    else:
          comp = category
          component_data, component_sheet, primary_col = reading_excel_sources(comp.strip(), sql_conn,cursor)
          valid_files = glob.glob(valid_path + '*.csv')
          #print(valid_files)
          for files in valid_files:            
            for sheet in component_sheet:
                head, tail = os.path.split(files)
                file_name = tail.rsplit('.',1)[0]    
                #print('ll',file_name,sheet) 
                if file_name.strip()==sheet.strip():
                  print(file_name,sheet)   
                  sheet_name = sheet.strip()
                 # print('primary_col',primary_col)
                  for primary in primary_col:
                   # print(primary)
                    reading_excel_data_from_source(valid_path, files, component_data, primary,comp,sql_conn,cursor,product_inscope_df, 
                                                   unstructure_processed_data_query,excel_date_found,category,sheet_name) 
                    
  except Exception as e:
    logger.error("Error in excel_extract2_key_value_pair function while processing {}".format(valid_path),exc_info=True)    
#***************************************************************************************************************************************
#function name: image_data_extract
#Objectiv: To extract required images from file
#input parameters:
#unstruct_data_key_info: will hold all the data except key_value extract data like(product_type, category, product)
#raw_df: will hold all the staging file path in dataframe which helps to move file to processed folder
#data_extract: will hold the key-value data 
#Usage: common code is written which extracts image structure for the identified product in the files based on the coordinates data #produced by tesseract ocr and update unstruct_data_key_info dataframe       
#called by: key_value_extract
#***************************************************************************************************************************************
def image_data_extract(file):
  try:        
      file=file.replace("dbfs:","/dbfs").strip()      
      logger.info("Executing image_Key_extract function")
      check_path=True
      temp_path = file.rsplit('staging',1)[0]  + 'temp/'
      print(temp_path,temp_path)
      intialize_temp_files(temp_path)
      #Conveting pdf to image file
      logger.info("Calling pdf to image conversion function")
      pdf_to_image_converison(file,temp_path)
      target= temp_path
      target_list= glob.glob(target+'*.*')
      #print('target_list',target_list)
      return target_list
  except Exception as e:
    logger.error("Error in image_data_extract function while processing {}".format(file),exc_info=True)
      
# To convert image to text with coordinates
def image_to_data_conversion(opened_image,file):
  try: 
    logger.info("Executing image_to_data_conversion function")
    txt=pyt.image_to_data(opened_image)
    txt_read=txt.split('\n')
    coordinates=[]
    last_y1=0
    sentence=''
    word_cords=[]
    line_cords=[]
    line_y1=0
    line_x1=0    
    for i in range(1,len(txt_read)):
      try:
          cords_str=str(txt_read[i]).split('\t')
          json={}
          cords={}
          text=cords_str[-1].strip()
          if len(cords_str)>10 and len(text)>0:
              word=cords_str[-1]
              json['text']=word
              y1=int(cords_str[7])
              cords["x1"]=cords_str[6]
              cords["y1"]=cords_str[7]
              cords["x2"]=int(cords_str[8])+int(cords_str[6])
              cords["y2"]=int(cords_str[9])+int(cords_str[7])
              json["coordinates"]=cords
              y_dif=(last_y1-y1)
              if y_dif <0:
                  y_dif=-1*y_dif
              if y_dif <=20  and last_y1>0:
                  sentence+=word+" "
                  word_cords.append(json)
              else:
                  if len(word_cords)>0:
                      line_json={"text":sentence.strip(),
                                 "coordinates":{"x1":line_x1,"y1":line_y1,
                                                "x2":word_cords[-1]["coordinates"]["x2"],
                                                "y2":word_cords[0]["coordinates"]["y2"]}}
                      line_cords.append(line_json)
                      final={"line_cords":line_cords,
                             "word_cords":word_cords}
                      word_cords=[]
                      line_cords=[]
                      coordinates.append(final)
                  line_x1=int(cords_str[6])
                  line_y1=int(cords_str[7])           
                  sentence=''
                  sentence+=word+" "
                  last_y1=y1
                  word_cords.append(json)
      except Exception as e:
        logger.error("Error in image to data conversion: inner iteration",exc_info=True)
    line_json={"text":sentence, 
               "coordinates":{"x1":line_x1,
                              "y1":line_y1,
                              "x2":word_cords[-1]["coordinates"]["x2"],
                              "y2":word_cords[-1]["coordinates"]["y2"]}}
    line_cords.append(line_json)
    final={"line_cords":line_cords,
           "word_cords":word_cords}
    coordinates.append(final)
    logger.info("Text data with their coordinates has been extracted successfully from image file "+file)
    return coordinates
  except Exception as e:
    logger.error("Error in image_to_data_conversion",exc_info=True)
    logger.error("Error in image file",file)
    
def image_to_cordinates(sql_conn,cursor,img_path,product_type_list,product_list,file_loc,category,file_name,img_count,unstructure_processed_data_query):
  block_json = {}
  json_check = []
  try:
      im = Image.open(img_path) 
      width, height=im.size
      #Converting image to text with coordinates as json doc
      logger.info("Calling image to data conversion function")
      coordinates=image_to_data_conversion(im,img_path)
      #print('img_pathkkkkkkkkkkk',img_path)
      crop_json={}
      crop_details=[]
      first_count=0
      expected_left=0
      diff_x1=0
      match_data_filter =[]
      for item in coordinates:
        try:
          line=item['line_cords']
          for ele in line:
              text=ele["text"].strip()
             # print(text)
              x1=int(ele["coordinates"]["x1"])
              y1=int(ele["coordinates"]["y1"])
              x2=int(ele["coordinates"]["x2"])
              y2=int(ele["coordinates"]["y2"])
             # print('text',text)
              match_f = None
              for prod_name in product_list:  
                #print(prod_name)
                rgx_img = re.compile('{}'.format(prod_name.replace('*','\*')),re.I)
                for match in re.finditer(rgx_img,text):
                  match_f = prod_name.strip()
                 # print('kamal',match_f)
                  break
              if  match_f is not None:# and 
                  if match_f not in match_data_filter  :    
                      match_data_filter.append(match_f)
                      first_count+=1
                      if first_count==1:
                          expected_left=x1
                      left=x1
                      top=y1-5
                      if(crop_json):
                          crop_json["bottom"]=y1-5
                          crop_details.append(crop_json) 
                      crop_json={}
                      crop_json={"left":0,
                                "top":top,
                                "right":width,
                                "name":match_f}
                  else:
                        crop_json={"left":0,"top":0,"right":width,"bottom" : height,"name":match_f}
                        crop_details.append(crop_json)
              elif(len(text)>1 and first_count>1):
                  diff_x1=expected_left-x1
                  if diff_x1<0:
                      diff_x1=diff_x1*-1
                  if (diff_x1<=3):
                      crop_json["bottom"]=y1-5
                      crop_details.append(crop_json)
                      crop_json={}                        
              if first_count==1:
                  first_count+=1
        except Exception as e:
          logger.error("Error in extracting image_to_cordinates function : coordinates iteration",exc_info=True)
      if(crop_json):
          if (y2-y1)>50:
              crop_json["bottom"]=y2+5
          else:
              crop_json["bottom"]=height
          crop_details.append(crop_json)
      else:
        crop_json={"left":0,"top":0,"right":width,"bottom" : height,"name":"file_name"}
        crop_details.append(crop_json)
      #item_image_name = []
      if len(crop_details)>0:        
        for item in crop_details:
            im1 = im.crop((item.get('left'), item.get('top'), item.get('right'), item.get('bottom'))) 
            name=item.get('name')            
            name_db = name
            if name == 'file_name':
              img_path = file_loc + 'non-relavent/' 
              img_count = img_count + 1
              name = file_name + '_' + str(datetime.now())[:19]
              product_type = 'null'
              name_db = 'null'
            else:              
              prod_index = product_list.index(name)
              product_type = product_type_list[prod_index]
              img_path = file_loc + 'relavent/' 
            if not os.path.exists(img_path) :
              path_exists(img_path)
              check_path=False                 
            path=img_path+name+".png"
            if os.path.exists(path):
               path = path.rsplit(".png",1)[0] + '_' + str(datetime.now())[:19]  +".png"
            im1.save(path)
            block_json['file_path'] = path
            data_extract=json.dumps(block_json,ensure_ascii=False)
            if product_type == 'null':
              unstructure_processed_data(unstructure_processed_data_query,category,product_type,name_db,data_extract,0,sql_conn,cursor)
            else:
              unstructure_processed_data(unstructure_processed_data_query,category,product_type,name_db,data_extract,1,sql_conn,cursor)
      #print(img_count)
      return img_count
  except Exception as e:
    logger.error("Error in  image_to_cordinates function",exc_info=True)
    
#*********************************************************************************************************************************
#function name: tex_Key_extract
#Ojective: Key value extart for FDA files
#input parameters:
#unstruct_data_key_info: will hold all the data except key_value extract data like(product_type, category, product)
#raw_df: will hold all the staging file path in dataframe which helps to move file to processed folder
#data_extract: will hold the key-value data 
#Usage: common code is written which extracts required key value data based on the keywords for US-FDA categories and update the 
#       unstruct_data_key_info dataframe
#called by: key_value_extract
#*********************************************************************************************************************************** 
def text_Key_extract(file,filter_df,content,staging_raw_file_path,file_images):
  try: 
    file=file.replace("dbfs:","/dbfs").strip()
    file_name =  file.split('/')[-1].rsplit('.',1)[0]
    block_json = {}
    #******************
    #value_extract
    #******************
    for  index_df in filter_df.index:
        start_string_index = None
        end_string_index = None
        rgx = re.compile(r'({})'.format(filter_df['start_key'][index_df]),re.I)
        #******************************
        #checking index of start key
        #******************************
        for match in re.finditer(rgx,content):
            if match.group():
                start_string_index = match.start()
                break

        if  start_string_index is not None:
            
            #******************************
            #checking index of end key
            #******************************
            if filter_df['end_key'][index_df].lower().strip().startswith('regex'):
              end_rex=filter_df['end_key'][index_df].lower().strip().split('regex:')[1]                   
              rgx = re.compile(r'({})'.format(filter_df['start_key'][index_df]),re.I)
              #print(end_rex)
              rgx_flag =''  
              #print('regex')
              #print(start_string_index)
              #******************************
              #checking index of start key
              #******************************              
              for match in re.finditer(rgx,content):
                  if match.group():
                      start_string_index = match.start()
                      rgx = re.compile(r'({})'.format(end_rex),re.I)
                      for match in re.finditer(rgx,content[start_string_index:]):
                        if match.group():
                            end_string_index = start_string_index + match.end() 
                            print('end_string_index',end_string_index)
                            rgx_flag = 's'
                            break 
                      if rgx_flag =='s':
                          break
                                                  
            elif filter_df['end_key'][index_df].lower().strip() == 'image':              
                #print('staging_raw_file_path',staging_raw_file_path.values[0])
                target_temp=image_data_extract(staging_raw_file_path)
                #print('file_images',file_images)
                if not os.path.exists(file_images):
                  path_exists(file_images)
                image_flag = ''
                for img_path in range(0,len(target_temp)):
                  im = PIL.Image.open(target_temp[img_path])
                  if im.mode=='P':
                    im = im.convert(palette=0)
                  im1 = im.filter(ImageFilter.EDGE_ENHANCE_MORE)                                    
                  config1 = (' --psm 6')
                  image_text = pyt.image_to_string(im1, config=config1)                          
                  rgx = re.compile(r'({})'.format(filter_df['start_key'][index_df]),re.I)                    
                  for match in re.finditer(rgx,image_text):
                    if match.group():                      
                      dbutils.fs.cp(target_temp[img_path].replace("/dbfs","dbfs:").replace('//','/')
                      ,file_images.replace("/dbfs","dbfs:").replace('//','/'))
                      block_json['image_path']  = file_images +  target_temp[img_path].rsplit('/',1)[1]
                      image_flag='s'
                      break
                  if image_flag == 's':
                    break   
  
            else:  
                rgx = re.compile(r'({})'.format(filter_df['end_key'][index_df]),re.I)
                for match in re.finditer(rgx,content[start_string_index:]):
                  if match.group():
                      end_string_index = start_string_index + match.end()  
                      break                                                                                                               
                      
        if  start_string_index is not None and end_string_index is not  None :#and end:
            text_extract = content[start_string_index:end_string_index].replace('\n',' ')#.replace('\u2014'
        
            #********************************************************************************
            #replacing end_key text in extracted text if it is presnt in start key column
            #********************************************************************************
            if not filter_df[filter_df['start_key'].str.contains(filter_df['end_key'][index_df])].empty:
                find_replace=text_extract.lower().find(filter_df['end_key'][index_df].lower())
                text_extract = text_extract[:find_replace]

            #**********************************************************************************
            #Writing the extracted text in the json based on the field names present in the 
            #***********************************************************************************
            
            if pd.isnull(filter_df['field'][index_df]):
                if 'data' not in block_json.keys():
                    block_json['data'] = text_extract
                else:
                    block_json['data'] = block_json['data'] + ' ' + text_extract
            else:
                if filter_df['field'][index_df] not in block_json.keys():
                    block_json[filter_df['field'][index_df]] = text_extract
                else:
                    block_json[filter_df['field'][index_df]] = block_json[filter_df['field'][index_df]] + ' ' + text_extract

    #if bool(block_json):
    #****************
    #date extract
    #****************
    rgx_pattern_date =["\d{1,2}\s*\/\d{1,2}\s*\/\d{4}","[a-zA-Z]{3,11}\s*\d{1,2}\s*,\s*\d{4}","\d{1,2}\s[a-zA-Z-]{3,11}\s*\d{4}",
                       "\d{1,2}\s*\-\s*[a-zA-Z]{3,11}\s*\-\s*\d{4}"] 

    date_result=re.search("|".join(rgx_pattern_date), content)   
    print(block_json)
    if 'date' not in block_json and 'Date' not in block_json:      
      if date_result is not None:
         block_json['Date'] = date_result.group().strip()
      else:
        block_json['Date'] = 'null'

    #********************
    # subject extract:
    #********************
    rgx_pat = r'\nre(:|\s)'
    sub_first = re.finditer(rgx_pat, content)
    for m_string in sub_first:
        sub_first_check_v1 = m_string.start(0)
        break
    sub_first_check_v2 = content.find('eu food contact statement')
    sub_last = content.find('\ndear')
    sub_first_check_v3 = content.rfind('\n',0,sub_last)    
    if sub_last != -1 and sub_first_check_v1 != -1 and sub_last > sub_first_check_v1 :
        subject = content[sub_first_check_v1:sub_last].strip()
    elif sub_last != -1 and sub_first_check_v2 != -1 and sub_last > sub_first_check_v2 :
        subject = content[sub_first_check_v2:sub_last].strip()
    elif sub_last != -1 and sub_first_check_v3 != -1 and sub_last > sub_first_check_v3 :
        subject = content[sub_first_check_v3:sub_last].strip()
    else:
        subject = file.split('/')[-1][:-4]
        block_json['subject'] = subject
   
    logger.info("data extract is successful for this {}".format(file) ) 
#     else:
#       logger.info("data extract is not found for this {}".format(file) )
#    print(block_json)
    return block_json
 
  except Exception as e:
    logger.error("Error in text_Key_extract function while processing {}".format(file),exc_info=True)
    
    
#**************************************************************************************************************************************
#function name: relavent_image_extract
#Ojective: TO identify relevant and non-relevant files
#input parameter: 
#file: will hold the absolute file path of all-text folder
#file_loc: will hold the valid folder file path 
#bdt_list: will hold all the BDT data fetched from product_inscope_df
#nam_prod_list: will hold all the NAM PROD data fetched from product_inscope_df
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#category_list: will store the category type in which extracted data falls into(like BDT, NAM PROD)
#product_type_list: will store product type which identified in extracted data
#file_path_list: will store absolute file path of all-text folder
#Usage: To differentiate relevant and non-relavent files based on the product inscope details like(NAM PROD, BDT, SILICONE US-FDA and EU
#file_validation: Moves the file to relevant and non-relevant folder based on the product_inscope
#called by: pattern_match_validation
#**************************************************************************************************************************************     
def relavent_image_extract(file,file_loc,content,product_inscope_df,category,file_is_valid_query,file_unique_list,sql_conn,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query):  
  global image_folder_list  
  try:             
      prod_flag =''
      product_list = []
      product_list_lower= []
      product_type_list =[]
      data_extract={}
      target_temp=image_data_extract(file)    
      file_name =  file.split('/')[-1].rsplit('.',1)[0]
      img_count = 0
     # print('file_name',file_name)  
      #*********************************************************
      #checking the  PROD present in the exrtacted content 
      #********************************************************* 
                                                 
      for prod_index in product_inscope_df.index:
        if not pd.isnull(product_inscope_df['Type'][prod_index]):
           try:
              prod_rgx = re.search(r'(([^a-zA-Z]|^){}[^a-zA-Z])'.format(re.escape(product_inscope_df['Text']
                                                                                  [prod_index].strip())),content,re.I)   
              if(prod_rgx):
                  prod_txt =  product_inscope_df['Text'][prod_index].strip().upper()
                  prod_flag = 's'                
                  logger.info('{} Successfully passed the inscope validation by containing {} {} in the \
                              content'.format(file,product_inscope_df['Type'][prod_index].strip(),prod_txt)) 
                  if prod_txt.strip() not in product_list:
                    #product_list_lower.append(prod_txt.lower().strip())
                    product_list.append(prod_txt.strip())
                    product_type_list.append(product_inscope_df['Type'][prod_index].strip())
                  
           except Exception as e:
                logger.error("Error in relavent_image_extract function inner iteration",exc_info=True)
      print('len(target_temp)',len(target_temp))
      print(file_name)
      if  'sip'  in file_name.lower():
        if len(target_temp) > 2:
          #print(target_temp)
          target_temp = target_temp[:2]
      #print(target_temp)
      for img_path in target_temp:        
        img_count = image_to_cordinates(sql_conn,cursor,img_path,product_type_list,product_list,file_loc,category,
                                        file_name,img_count,unstructure_processed_data_query)

               
  except Exception as e:
      logger.error("Error in relavent_image_extract function: outer iteration",exc_info=True)
          
def valid_files_copy(file,file_valid_type,data_extract):    
  try:
    if not os.path.exists(file_valid_type):
      path_exists(file_valid_type)
    text_name = file_valid_type.replace("dbfs:","/dbfs") + file.split('/')[-1][:-4] + '.txt'
    with open(text_name, "w",encoding='utf8') as file_write:
       json.dump(data_extract,file_write,ensure_ascii=False)
    logger.info('key-data extract of file {} has been written into {}'.format(file,text_name))    
  except Exception as e:
    logger.error('something went wrong in valid_files_copy functionwhile processing \
                         {}'.format(file),exc_info=True)
      
#**************************************************************************************************************************************
#function name: relavent_text_extract
#Ojective: TO identify relevant and non-relevant files
#input parameter: 
#file: will hold the absolute file path of all-text folder
#file_loc: will hold the valid folder file path 
#bdt_list: will hold all the BDT data fetched from product_inscope_df
#nam_prod_list: will hold all the NAM PROD data fetched from product_inscope_df
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#category_list: will store the category type in which extracted data falls into(like BDT, NAM PROD)
#product_type_list: will store product type which identified in extracted data
#file_path_list: will store absolute file path of all-text folder
#Usage: To differentiate relevant and non-relavent files based on the product inscope details like(NAM PROD, BDT, SILICONE US-FDA and EU
#file_validation: Moves the file to relevant and non-relevant folder based on the product_inscope
#called by: pattern_match_validation
#**************************************************************************************************************************************   
def relavent_text_extract(file,staging_file,file_loc,content,product_inscope_df,category,file_is_valid_query,file_unique_list,sql_conn,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query,sil_elast_product_list= None):
  try: 
      file_relavent = file_loc + 'relavent/'
      file_non_relavent = file_loc + 'non-relavent/'
      prod_flag =''      
      sil_prod_flag =''
      data_extract={}
      filter_df= unstruct_category_key_df[unstruct_category_key_df['category'].str.contains(category)]
      file_name =  file.split('/')[-1].rsplit('.',1)[0].strip()
      file_images = file_loc + 'relavent_images/' + file_name.replace('*','') + '/'   
#       external_staging_path_rep = external_processed_files_df_rep['blob_staging_file_path'].values.tolist()
#       external_satging_file_name_rep = external_processed_files_df_rep['file_name'].values.tolist()
#       external_satging_file_format_rep = external_processed_files_df_rep['file_format'].values.tolist()
      staging_raw_file_path = staging_file.strip()
#       try:
#         staging_raw_file_path = raw_df['file_name'].astype('str').str.contains(file_name.replace('*','\*'))
#         #staging_raw_file_path = raw_df['file_name'].astype('str').str.contains(file_name)
#         staging_raw_file_path=raw_df[staging_raw_file_path]['file_name'].values.tolist()[0]
#       except Exception as e:
#         if file_name in  external_satging_file_name_rep:
#           staging_file_index= external_satging_file_name_rep.index(file_name)
#           staging_raw_file_path= external_staging_path_rep[staging_file_index] + file_name +  
#           external_satging_file_format_rep[staging_file_index] 
#           print('staging_raw_file_path1',staging_raw_file_path)
      processed_path = file.rsplit('analytics',1)[0]  + 'Processed/'      
      if not filter_df.empty:
        data_extract=text_Key_extract(file,filter_df,content,staging_raw_file_path,file_images)    
      if data_extract is None:
        data_extract={}
#       if 'image_path' not in data_extract:
#         if not os.path.exists(file_loc + 'non-relavent/')
#         dbutils.fs.cp(data_extract_copy['image_path'].replace("/dbfs","dbfs:"),file_relavent)
#         logger.info('{} it falls under incsope product, So moving this to {}'.format(file,file_relavent))     
 
     #*****************************************
     #Copying files to processed folder
     #*****************************************
      #if bool(data_extract):
    #  print('staging_raw_file_path',staging_raw_file_path)
     # print(staging_raw_file_path.shape)
      #print('staging_raw_file_path.values',staging_raw_file_path.values)        
      if not os.path.exists(processed_path):        
        path_exists(processed_path)
      dbutils.fs.cp(staging_raw_file_path.replace("/dbfs","dbfs:").replace('//','/').strip(), 
                    processed_path.replace("/dbfs","dbfs:"))
      #print('nnnnn',processed_path.replace("dbfs:","/dbfs"),'mm',staging_raw_file_path.values[0].rsplit('/',1)[1])
      data_extract['file_path']  = processed_path.replace("dbfs:","/dbfs") + staging_raw_file_path.rsplit('/',1)[1]
      data_extract['file_name']  = staging_raw_file_path.rsplit('/',1)[1]
      data_extract_copy=data_extract
      data_extract = json.dumps(data_extract) 
      #*********************************************************
      #checking the PROD present in the exrtacted content 
      #*****************************************************prod_index**** 
      for prod_index in product_inscope_df.index:        
        if not pd.isnull(product_inscope_df['Type'][prod_index]):
           try:
              if not str(product_inscope_df['Text'][prod_index].strip()).isspace():
                prod_rgx = re.search(r'(([^a-zA-Z]|^){}[^a-zA-Z])'.format(re.escape(product_inscope_df['Text']
                                                                                    [prod_index].strip())),content,re.I)   
                if(prod_rgx):
                    prod_txt =  product_inscope_df['Text'][prod_index].strip().upper()
                    prod_flag = 's'                
                    logger.info('{} Successfully passed the inscope validation by containing {} {} in the \
                                content'.format(file,product_inscope_df['Type'][prod_index].strip(),prod_txt))   
                    unstructure_processed_data(unstructure_processed_data_query,category,product_inscope_df['Type']
                                               [prod_index].strip(),prod_txt,data_extract,1,sql_conn,cursor)
           except Exception as e:
                logger.error("Error in relavent_file function inner iteration while processing {}".format(file),exc_info=True)  

      #*************************************************************************************************************
      #checking the EU and US-FDA present in the exrtacted content if NAM PROD and BDT not in the content
      #*************************************************************************************************************
      if prod_flag != 's'  and sil_elast_product_list != None:            
          try:
            for sil_prod in sil_elast_product_list:
                sil_prod_rgx = re.search(r'(([^a-zA-Z]|^){}[^a-zA-Z])'.format(re.escape(sil_prod)),content,re.I)          
                if(sil_prod_rgx):
                    sil_prod_flag = 's'
                    sil_prod = sil_prod.upper()
                    logger.info('{} Successfully passed the inscope validation by containing silicone elatomer  {} in the \
                                content'.format(file,sil_prod))
                    unstructure_processed_data(unstructure_processed_data_query,category,'Silicone\
                                 elastomer',sil_prod,data_extract,1,sql_conn,cursor)
          except Exception as e:
            logger.error("Error in relavent_file function: silicone_elast_prod inner iteration while processing \
                         {}".format(file),exc_info=True)

      #*************************************************************************************************************
      #Moving the files to Non-relevant folder if NAM PROD, BDT, EU and US-FDA not in the content
      #*************************************************************************************************************
      if prod_flag != 's'   and sil_prod_flag != 's' :         
            logger.info('{} it does not fall under incsope product, So moving this to {}'.format(file,file_non_relavent))
            valid_files_copy(file,file_non_relavent,data_extract) 
            unstructure_processed_data(unstructure_processed_data_query,category,'null','null',data_extract,0,sql_conn,cursor)
      else: 
           # if 'image_path' not in data_extract_copy:              
              valid_files_copy(file,file_relavent,data_extract)
              logger.info('{} it falls under incsope product, So moving this to {}'.format(file,file_relavent))             
            
  except Exception as e:
      logger.error("Error in relavent_file function: outer iteration",exc_info=True)
          
          
#******************************************************************************************************************************************
#function name : pattern_match_validation
#Ojective: Pattern match validation on each file extracted from staging path
#input parameters:
#external_processed_files: will have all the extracted file paths from the file_processing_info table
#external_staging_path: will have all the staging folder path from the file_processing_info table for the extracted file
#external_satging_file_format: will have all file formats from the file_processing_info table for the extracted file
#pattern_matching_query: will have select query for the pattern_matching_keys table
#pattern_key_df: will store all pattern matching keywords by passing pattern_matching_query to external_source_data function
#product_inscope_query: will have the select query of product_inscope(Nam prod, BDT,CAS No details)
#product_inscope_df: will hold all product inscope details in dataframe by passing product_inscope_query to external_source_data function
#bdt_list: will hold all the BDT data fetched from product_inscope_df
#nam_prod_list: will hold all the NAM PROD data fetched from product_inscope_df
#silicone_elastomer_product_query: will have select query for the silicone elastomer table which contains 
#                                  valid US-FDA and EU Product data extracted from Silicone elastomer brochure pdf file
#                                  scrapped from momentive website
#silicone_elastomer_product_df: will hold all the US-FDA and EU in dataframe by passing product_inscope_query to external_source_data 
#                               function
#relavent_file_extract: To differentiate relevant and non-relavent files based on the product inscope details
#Usage: This functionality helps to perform 15 pattern match validation on each extracted text files based on all_text path stored in #file_processing_info table
#output parameter:
#unstruct_data_df: will hold all the valid file paths, categories and product keys which will be input for key value extract function
#******************************************************************************************************************************************   
def pattern_match_validation(sql_conn,external_processed_files_df,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query):
  try:
    logger.info('Executing pattern_match_validation for all exctracted files')
    external_processed_files = external_processed_files_df['blob_all_txt_file_path'].values.tolist()
    external_excel_files = external_processed_files_df['excel_valid_path'].values.tolist()
    external_staging_path = external_processed_files_df['blob_staging_file_path'].values.tolist()
    external_satging_file_name = external_processed_files_df['file_name'].values.tolist()
    external_satging_file_format = external_processed_files_df['file_format'].values.tolist()
#     external_staging_path_strip_list = []
#     external_satging_file_name_strip_list = []
#     external_satging_file_format_strip_list = []
#     for index_in in range(len(external_staging_path)):
#       external_staging_path_strip_list.append(external_staging_path[index_in].strip())
#       external_satging_file_name_strip_list.append(external_satging_file_name[index_in].strip())
#       external_satging_file_format_strip_list.append(external_satging_file_format[index_in].strip())
#     external_processed_files_df_rep = pd.DataFrame()
#     external_processed_files_df_rep['blob_staging_file_path'] = external_staging_path_strip_list
#     external_processed_files_df_rep['file_name'] = external_satging_file_name_strip_list
#     external_processed_files_df_rep['file_format'] = external_satging_file_format_strip_list
    pattern_matching_query = config.get('mount_path','pattern_match')  
    pattern_key_df = external_source_data(sql_conn,pattern_matching_query)  
    pattern_category = list(set(pattern_key_df['pattern_category'].values.tolist()))
    pattern_image_list = pattern_key_df[pattern_key_df['result_type'].astype('str').str.contains('1')]['pattern_keys'].values.tolist()
    file_is_valid_query = config.get('mount_path', 'file_is_valid')
    product_inscope_query = config.get('mount_path','product_inscope')
    product_inscope_df = external_source_data(sql_conn,product_inscope_query)
    product_inscope_df.replace(r'^\s*$', np.nan, regex=True,inplace=True)
    product_inscope_df = product_inscope_df.fillna('null')
    product_inscope_df = product_inscope_df.fillna('null')
   # print(product_inscope_df.columns)
    silicone_elastomer_product_query = config.get('mount_path','silicone_elastomer_product')
    silicone_elastomer_product_df =  external_source_data(sql_conn,silicone_elastomer_product_query) 
    matnbr_list = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('MATNBR')]['Text1'].values.tolist()))
    bdt_list = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('MATNBR')]['Text3'].values.tolist()))
    nam_prod_list = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('NAMPROD')]['Text1'].values.tolist()))
    nam_prod_list_FDA = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('NAMPROD') & 
                                                    product_inscope_df['SUBCT'].str.contains('REAL_SUB')]['Text1'].values.tolist()))
    cas_list = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('NUMCAS')]['Text1'].values.tolist()))
    product_type_list_fda = ['NAMPROD'] * len(nam_prod_list_FDA) + ['BDT'] * len(bdt_list) +  ['NUMCAS'] * len(cas_list)
    product_type_list= ['NAMPROD'] * len(nam_prod_list) + ['BDT'] * len(bdt_list) +  ['NUMCAS'] * len(cas_list) + ['MATNBR'] * \
    len(matnbr_list)
    product_valid_list = nam_prod_list + bdt_list + cas_list + matnbr_list
    product_valid_list_fda = nam_prod_list_FDA + bdt_list + cas_list
    product_inscope_df = pd.DataFrame(columns=['Type', 'Text'])
    product_inscope_df['Type'] = product_type_list
    product_inscope_df['Text'] = product_valid_list
    product_inscope_df_fda = pd.DataFrame(columns=['Type', 'Text'])
    product_inscope_df_fda['Type'] = product_type_list_fda
    product_inscope_df_fda['Text'] = product_valid_list_fda
    silicone_elastomer_product_query = config.get('mount_path','silicone_elastomer_product')
    silicone_elastomer_product_df =  external_source_data(sql_conn,silicone_elastomer_product_query)  
    silicone_elastomer_product_df=silicone_elastomer_product_df.rename(columns = {'eu_fda':'EU-FDA','us_fda':'US-FDA'})
    file_unique_list =[]
      
    #**********************************************
    #Iterating each files for pattern matching 
    #**********************************************
    for index in range(0,len(external_processed_files)):
      try:
        image_falg =''
        analytics_valid_path = external_processed_files[index].rsplit('all-text',1)[0] + 'valid-files/'
        analytics_invalid_path = external_processed_files[index].rsplit('all-text',1)[0] + 'invalid-files/'
        if analytics_invalid_path not in valid_folder_list:
          #path_exists(analytics_valid_path)
          #path_exists(analytics_invalid_path)
          valid_folder_list.append(analytics_invalid_path)        
        file=external_processed_files[index].replace("dbfs:","/dbfs")     
        if file in extracted_file_list:
          content = open(file, 'r', encoding = 'utf-8').read()
          file_valid_flag =''

          #****************************************************************************************************
          #checking pattern_match on each file based on pattern_category,pattern_keys and filter_condition
          #***************************************************************************************************
          for pattern_cat_match in pattern_category:
              pattern_match_flag = '' 
              #*************************************
              #filtering based on pattern_category
              #*************************************
              pattern_filter_condition_df = pattern_key_df[pattern_key_df['pattern_category'].str.
                                            contains("^\s*{}\s*$".format(pattern_cat_match),case=False)] 
              and_condition_list = []
              or_condition_list = []
              if not pattern_filter_condition_df.empty:
                  #**********************************************************
                  #filtering based on pattern_keys and filter_condition
                  #**********************************************************
                  and_condition_df = pattern_filter_condition_df[pattern_filter_condition_df['filter_condition'].str.
                                                                 contains("^\s*{}\s*$".format('1'),case=False)]
                  or_condition_df = pattern_filter_condition_df[pattern_filter_condition_df['filter_condition'].str.
                                                                contains("^\s*{}\s*$".format('0'),case=False)]
                                    
                               
                      
                  #*****************************'*********************************************
                  #checking if both filter condition '1' or '0' present in the category
                  #**************************************************************************
                  if  all([and_condition_df.empty,or_condition_df.empty]) == True:                 
                      and_condition_list = and_condition_df['pattern_keys'].values.tolist()                                        
                      or_condition_list = or_condition_df['pattern_keys'].values.tolist()

                      if all(match.lower().strip() in content.lower() for match in and_condition_list):
                          if any(match.lower() in content.lower() for match in or_condition_list):
                              pattern_match_flag = 's'
                  #******************************************
                  #checking only filter condition 'and' 
                  #******************************************
                  elif not and_condition_df.empty:
                      and_condition_list = and_condition_df['pattern_keys'].values.tolist()
                      if all(match.lower().strip() in content.lower() for match in and_condition_list):
                          pattern_match_flag = 's'
                  #******************************************
                  #checking only filter condition 'or' 
                  #******************************************
                  elif not or_condition_df.empty:
                      or_condition_list = or_condition_df['pattern_keys'].values.tolist()
                      if any(match.lower().strip() in content.lower() for match in or_condition_list):
                          pattern_match_flag = 's'
              #****************************************************
              #if the file under goes any of the pattern category
              #****************************************************                             
              if  pattern_match_flag == 's':      
                  logger.info('{} is found in {}'.format(file,pattern_cat_match))
                  #*****************************************************************************************
                  #validating for relavent and non-relavant file if the key-value extract is text from file
                  #*****************************************************************************************
                  print('pattern_cat_match',pattern_cat_match)
                  pattern_result_list =(pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$"
                         .format(pattern_cat_match),case=False)])['result_type'].values.tolist() 
                  print('pattern_result_list',pattern_result_list)
                  if all(int(match.lower().strip()) == 0 for match in pattern_result_list):                         
                      logger.info('{} its a text extraction type so moving this file to relavent_text_extract function'.format(file))
                      sil_elast_product_list = None
                      for match_mpm_cat in silicone_elastomer_product_df.columns:
                          if match_mpm_cat.lower().strip() == pattern_cat_match.strip().lower():                            
                              sil_elast_product_list = list(set(silicone_elastomer_product_df[match_mpm_cat].values.tolist())) 
                      staging_file = external_staging_path[index] + file.rsplit('/',1)[1].rsplit('.',1)[0] + \
                        external_satging_file_format[index]
                      if external_satging_file_format[index].strip()  in ['.pdf','.docx','.doc','.pptx']:
                        file_loc = file.rsplit('all-text',1)[0] + 'valid-files/'+ pattern_cat_match.strip() + '/'
                        file_is_valid = file_is_valid_query.format(1,1,'null',file.replace("dbfs:","/dbfs"))                      
                        update_operation(file_is_valid,sql_conn,cursor)
                        if pattern_cat_match.strip() in ('US_FDA','EU-FDA'):
                            product_inscope_df = product_inscope_df_fda
                        #print('acat',pattern_cat_match.strip())                                            
                        relavent_text_extract(file,staging_file,file_loc,content,product_inscope_df,pattern_cat_match.strip(),
                          file_is_valid_query,file_unique_list,sql_conn,cursor,unstruct_category_key_df,                      
                          raw_df,unstructure_processed_data_query,sil_elast_product_list)
                                 
                        file_valid_flag ='s'

                  #******************************************************************************************
                  #validating for relavent and non-relavant file if the key-value extract is image from file
                  #****************************************************************************************** 
                 
                  elif all(int(match.lower().strip()) == 1 for match in 
                           (pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$"
                           .format(pattern_cat_match),case=False)])['result_type'].values.tolist()): 
                      and_or_list = and_condition_list + or_condition_list
                      #print(and_or_list)
                      image_file_name = file.rsplit('/',1)[1]
                      if any(match.lower().strip() in image_file_name.lower().strip() for match in and_or_list) or 'sip' in \
                        image_file_name.lower():
                        file_loc = file.rsplit('all-text',1)[0] + 'valid-files/'+ pattern_cat_match.strip() + '/'
                        #print('file_loc,',file_loc)
                        #*********************************************************************************
                        #file: will hold the file present in the statging path for process using tesseract
                        #*********************************************************************************
                        #print(file.replace("dbfs:","/dbfs"))
                        file1 = external_staging_path[index] + file.rsplit('/',1)[1].rsplit('.',1)[0] + \
                        external_satging_file_format[index]
                        #print(file)
                        if pattern_cat_match.strip() in ('man_flow_diagram'):
                            product_inscope_df = product_inscope_df_fda
                        product_inscope_df = product_inscope_df_fda
                        if file1.strip().endswith('.pdf'):# and not file.endswith('.xlsm') and not file.endswith('.csv'):
                          file_is_valid = file_is_valid_query.format(1,1,'null',file1.replace("dbfs:","/dbfs"))
                          update_operation(file_is_valid,sql_conn,cursor)
                          logger.info('{} its a image extraction type so moving this file to relavent_image_extract function \
                                      '.format(file1))
                          relavent_image_extract(file1,file_loc,content,product_inscope_df,pattern_cat_match.strip(),file_is_valid_query,
                          file_unique_list,sql_conn,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query) 
                          file_valid_flag ='s'
                                                                                                                                                 
                  #******************************************************************************************
                  #validating for relavent and non-relavant file if the key-value extract  from excel
                  #******************************************************************************************  
                  elif all(int(match.lower().strip()) == 2 for match in 
                           (pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$"
                           .format(pattern_cat_match),case=False)])['result_type'].values.tolist()):     
                      file2 = external_processed_files[index].replace("dbfs:","/dbfs")
                      file_is_valid = file_is_valid_query.format(1,1,'null',file2.replace("dbfs:","/dbfs"))                      
                      update_operation(file_is_valid,sql_conn,cursor)
                      valid_path = external_excel_files[index]
                      logger.info('{} its a Excel extraction type so moving this file to excel_extract2_key_value_pair function \
                                    '.format(valid_path))
                      excel_date_found = excel_date(content,file)
                      print('valid_path',valid_path)
                      excel_extract2_key_value_pair(valid_path,sql_conn,cursor,pattern_cat_match.strip(),product_inscope_df, 
                                           unstructure_processed_data_query,excel_date_found)
#                       key_data_extract_external_source(valid_path,sql_conn,cursor,pattern_cat_match.strip(),
#                       unstructure_processed_data_query,excel_date_found)
                      file_valid_flag ='s'  
  
                  #******************************************************************************************
                  #validating for relavent and non-relavant file if the key-value extract  from tables in pdf
                  #******************************************************************************************  
#                   elif all(int(match.lower().strip()) == 3 for match in 
#                            (pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$"
#                            .format(pattern_cat_match),case=False)])['result_type'].values.tolist()):     
#                       file2 = external_processed_files[index].replace("dbfs:","/dbfs")
#                       file_is_valid = file_is_valid_query.format(1,1,'null',file2.replace("dbfs:","/dbfs"))                      
#                       update_operation(file_is_valid,sql_conn,cursor)
#                       valid_path = external_excel_files[index]
#                       logger.info('{} its a Excel extraction type so moving this file to excel_extract2_key_value_pair function \
#                                     '.format(valid_path))
#                       excel_date_found = excel_date(content,file)
#                       print('valid_path',valid_path)
#                       excel_extract2_key_value_pair(valid_path,sql_conn,cursor,pattern_cat_match.strip(),product_inscope_df, 
#                                            unstructure_processed_data_query,excel_date_found)
# #                       key_data_extract_external_source(valid_path,sql_conn,cursor,pattern_cat_match.strip(),
# #                       unstructure_processed_data_query,excel_date_found)
#                       file_valid_flag ='s' 
                       
          image_file_name = file.rsplit('/',1)[1]    
          #print(image_file_name)
          #print(pattern_image_list)
          if file_valid_flag !='s':            
              for match in pattern_image_list:                
                if match.lower().strip() in image_file_name.lower().strip(): 
                  #pattern_key_df[pattern_key_df['result_type'].astype('str').str.contains('1')]['pattern_keys'].values.tolist()
                  pattern_catch_found = pattern_key_df[pattern_key_df['pattern_keys'].astype('str').str.contains(match.strip())]
                  if pattern_catch_found.shape[0] == 1:
                    pattern_catch_found = pattern_catch_found.values.tolist()[0][1]
                    print(pattern_catch_found)
                  else:
                    pattern_catch_found = pattern_catch_found[pattern_category].values.list()[0]
                  file_loc = file.rsplit('all-text',1)[0] + 'valid-files/'+ pattern_catch_found.strip() + '/'
                  #*********************************************************************************
                  #file: will hold the file present in the statging path for process using tesseract
                  #*********************************************************************************
                  file1 = external_staging_path[index] + file.rsplit('/',1)[1].rsplit('.',1)[0] + \
                  external_satging_file_format[index]
                  if pattern_catch_found.strip() in ('man_flow_diagram'):
                     product_inscope_df = product_inscope_df_fda
                  #print(file)
                  if file1.strip().endswith('.pdf'):# and not file.endswith('.xlsm') and not file.endswith('.csv'):
                    file_is_valid = file_is_valid_query.format(1,1,'null',file1.replace("dbfs:","/dbfs"))
                    update_operation(file_is_valid,sql_conn,cursor)
                    logger.info('{} its a image extraction type so moving this file to relavent_image_extract function \
                                '.format(file1))
                    relavent_image_extract(file1,file_loc,content,product_inscope_df,pattern_catch_found.strip(),file_is_valid_query,
                    file_unique_list,sql_conn,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query) 
                    file_valid_flag ='s'
          #*************************************************************************************************
          #Moving the files to invalid-files folder as file content doesn't fall under pattern validation
          #*************************************************************************************************
          if file_valid_flag !='s':
              file_loc = file.rsplit('all-text',1)[0] + 'invalid-files/'
              if not os.path.exists(file_loc):
                path_exists(file_loc)
              file = file.replace("/dbfs","dbfs:")
              file_loc = file_loc.replace("/dbfs","dbfs:") 
              dbutils.fs.cp(file, file_loc) 
              file_name = file.rsplit('/',1)[-1]
              file_loc = file_loc.replace("dbfs:","/dbfs") + file_name
              file_is_valid = file_is_valid_query.format(0,0,file_loc,file.replace("dbfs:","/dbfs"))
              update_operation(file_is_valid,sql_conn,cursor)  
      except Exception as e:
        logger.error('file not found {}'.format(file),exc_info=True)
  except Exception as e:
    logger.error('something went wrong in pattern match validation',exc_info=True)
    
    
#**********************************************************************
#function name: copy_files
#Objective: To copy files from one folder to another
#input parameter:
#file_list: will hold all the files to moved from target in a list 
#staging_pdf_type: will hold the destination folder path
#called by: sharepoint_native_scanned_pdf_split
#**********************************************************************
def copy_files(file_list, staging_pdf_type):
    try:
        logger.info("Executing copy_files function")
        count=0
        if staging_pdf_type not in native_scanned_folder_list:
          path_exists(staging_pdf_type)
          native_scanned_folder_list.append(staging_pdf_type)
        for file in file_list:
          try:
              file=file.replace("/dbfs","dbfs:").replace("//","/")
              file_loc = staging_pdf_type.replace("/dbfs","dbfs:").replace("//","/")
              dbutils.fs.cp(file, file_loc)
              logger.info(file + ' copied to ' + staging_pdf_type)
              count+=1
          except Exception as e:
              logger.error("Error while copying ",file)
              logger.error("Error in copy_files function:iteration",exc_info=True)
        logger.info("Number of files copied to "+ staging_pdf_type+" : "+str(count))
    except Exception as e:
        logger.error("Error in copy_files function",exc_info=True)
#**************************************************************************************************************
#function name: excel2csv
#Ojective: excel to csv formats
#input parameters
#path: will hold the xlsx file path 
#Sheet: will hold sheet name present in the xlsx file
#Usage:converts excel file type into csv for text extarction as excel not supported in databricks
#called by: xlsx_text_extract
#**************************************************************************************************************
def excel2csv(path, sheet,staging_path):
  try:
      wb = openpyxl.load_workbook(path)
      sh = wb[sheet]
      head, tail = os.path.split(path)
      filename = path.rsplit('/')[-1].split('.')[0]
      file = head + '/' + 'temp/csv/' + sheet.strip() + '.csv'
      #path_exists(staging_path +'temp/csv/')
     # dbutils.fs.mkdirs((staging_path +'temp/csv/').replace("/dbfs","dbfs:").replace("//","/")) 
      with open(file, 'w', encoding="utf-8") as f:
          c = csv.writer(f)
          for r in sh.rows:
              c.writerow([cell.value for cell in r])
      return file
  except Exception as e:
    logger.error('Error in excel2csv function while converting {}'.format(path),exc_info=True)
    
#********************************************************************************************************************
#Function name: excel2txt
#Objective: Excel to text 
#input parameters
#staging_path : will hold the staging path of EXCEl type file fetched 
#abs_path: will hold file to be extracted 
#filename: will hold name of the file to be extracted
#Usage: common code is written which converts all the excel type files into text will be done using this function
#called by: csv_text_extract, xlsx_text_extract
#*********************************************************************************************************************
def excel2txt(staging_path, abs_path):
  try:
    data_text = pd.read_csv(abs_path, encoding='cp1252')
    file = staging_path +'temp/temp_all_text/' + abs_path.rsplit('/',1)[1].replace('.csv','.txt').strip()
    data_text.to_csv(file)    
  except Exception as e:
     logger.error('Error in excel2csv function while converting {}'.format(abs_path),exc_info=True)

#******************************************************************************************************************************
#function name: csv_text_extract
#Ojective: csv file into text extract
#input parameters
#staging_path : will hold the staging path of xlsx file fetched 
#csv_list: will hold all the csv files in a list
#source_type: will hold the respective source type of staging path
#all_files: will hold all-text path where text files to be stored which got extracted from the csv sheet
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#Usage: common code is written which converts all the csv file into text and stores the extracted data in all_files area in txt format,  #       then file path into the file_processing_info table
#called by : external_folder_structure_process
#****************************************************************************************************************************** 
def csv_text_extract(staging_path,csv_list,source_type,all_files,excel_files,file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor):
  try:
      logger.info('Executing csv_text_extract function') 
      for abs_path in csv_list:
        try:
          head, tail = os.path.split(abs_path)
          file_extn =tail.rsplit('.',1)[-1]
          file_name = tail.rsplit('.',1)[0]
          path_exists(staging_path+'temp/csv/')
          path_exists(staging_path +'temp/temp_all_text/')
          sheet =str(1)
          #**********************************************************************************************************************
          #excel2txt: It converts CSV file into text by taking  file path and sheet name present in the csv as input  
          #**********************************************************************************************************************
          excel2txt(staging_path, abs_path)
          text = glob.glob(staging_path +'temp/temp_all_text/'+'*.txt')
          #dbutils.fs.mkdirs(all_files.replace("/dbfs","dbfs:")) 
          text_csv = pd.DataFrame()
          file_path = all_files + file_name + '.txt'
          if not os.path.exists(excel_files + file_name +'/'):
            path_exists(excel_files + file_name +'/')
          dbutils.fs.cp(abs_path.replace("/dbfs","dbfs:"), (excel_files + file_name+'/').replace("/dbfs","dbfs:"), recurse=True) 
          for t in text:
            data = pd.read_csv(t, encoding='utf-8')
            text_csv = text_csv.append(data)
          text_csv.to_csv(file_path)
            #*******************************************************************************************************************
            #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
            #update_operation
            #*******************************************************************************************************************
          if file_path.replace('//','/').strip() not in file_processing_blob_all_txt_list:
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, \
            {},'{}')".format(source_type, file_name,'Excel','.csv', staging_path.replace('//','/'), file_path.replace('//','/').strip(),
                             1,0,'GETDATE()','GETDATE()',excel_files + file_name+'/')
          else:
              file_processing_info_query = update_file_processing_info.format('GETDATE()',file_path.replace('//','/').strip())
          update_operation(file_processing_info_query,sql_conn,cursor)
          logger.info('{}  extract_csv_text sucessfully'.format(file_path.replace('//','/')))
            
        except Exception as e:
          #********************************************************************************************************************
          #Creation of insert query for the extracted invalid file path to the file_processing_info table and executed using
          #update_operation
          #********************************************************************************************************************
          file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {},{},'{}')"\
          .format(source_type, file_name, 'Excel','.csv', staging_path.replace('//','/'), 'null', 0,0,'null','null','null')
          update_operation(file_processing_info_query,sql_conn,cursor)
          logger.error('Error in extracting csv_ text {}'.format(file_path.replace('//','/')),exc_info=True)
          
          
  except Exception as e:
    logger.error('Something went wrong in the csv_text_extract function', exc_info=True)  
    
#*****************************************************************************************************************************************
#function name : xlsx_text_extract
#objective : Extraction of text from excel sheets
#file_processing_info table
#input parameters
#staging_path : will hold the staging path of xlsx file fetched 
#xlsx_list: will hold all the xlsx files in a list
#source_type: will hold the respective source type of staging path
#all_files: will hold all-text path where text files to be stored which got extracted from the excel sheet
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#Usage: common code is written which converts all the xlsx file into text and stores the extracted data in all_files area in txt format, then file path into the file_processing_info table
#called by : external_folder_structure_process
#*****************************************************************************************************************************************      
def xlsx_text_extract(staging_path,xlsx_list,source_type,all_files,excel_files,file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor):
  try:
      for abs_path in xlsx_list:
        try:
          abs_path = abs_path.replace('//','/')
          staging_path = staging_path.replace('//','/')
          all_files =  all_files.replace('//','/')
          head, tail = os.path.split(abs_path)
          file_extn = tail.rsplit('.',1)[-1]
          file_name = tail.rsplit('.',1)[0]
          wb = openpyxl.load_workbook(abs_path) 
          allsheets = list(wb.sheetnames)
          path_exists(staging_path+'temp/csv/')
          path_exists(staging_path +'temp/temp_all_text/')
          file_path = all_files + file_name + '.txt'
          if not os.path.exists(excel_files + file_name +'/'):
            path_exists(excel_files + file_name +'/')         
          for sheet in allsheets:
            excel2csv(abs_path, sheet,staging_path)
          temp_path = glob.glob(staging_path+'temp/csv/'+'*.*')
          path_exists(staging_path +'temp/temp_all_text/')                    
          for i in range(len(temp_path)):
             # print(temp_path)
              dbutils.fs.cp(temp_path[i].replace("/dbfs","dbfs:"), (excel_files + file_name+'/').replace("/dbfs","dbfs:"), recurse=True)
              excel2txt(staging_path, temp_path[i])
          text_excel = glob.glob(staging_path +'temp/temp_all_text/'+'*.txt')
          text1 = pd.DataFrame()
          for t in text_excel:
            data = pd.read_csv(t, encoding='utf-8')
            text1 = text1.append(data)
          text1.to_csv(file_path)
              
          #**************************************************************************************************************
          #Creation of insert query for the extracted valid file path to the file_processing_info table using
          #update_operation
          #*************************************************************************************************************** 
          if file_path.replace('//','/').strip() not in file_processing_blob_all_txt_list:
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, \
            {},'{}')".format(source_type, file_name,'Excel','.csv', staging_path.replace('//','/'), file_path.replace('//','/').strip(),
                             1,0,'GETDATE()','GETDATE()',excel_files + file_name+'/')
          else:
              file_processing_info_query = update_file_processing_info.format('GETDATE()',file_path.replace('//','/').strip())
          update_operation(file_processing_info_query,sql_conn,cursor)
          logger.info('{}  extract_csv_text sucessfully'.format(file_path.replace('//','/')))          
                    
        except Exception as e:     
            #**************************************************************************************************************
            #Creation of insert query for the extracted invalid file path to the file_processing_info table using
            #update_operation
            #***************************************************************************************************************
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, \
            {},'{}')".format(source_type, file_name, 'Excel','.csv', staging_path.replace('//','/'), 'null', 0,0,'null','null','null')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.error('{}  is not extracted'.format(file_path.replace('//','/')))
            logger.error('Error in xlsx_text_extract function while converting {}'.format(abs_path),exc_info=True)
          
  except Exception as e:
    logger.error('Something went wrong in the xlsx_text_extract function', exc_info=True)       

def eml_text(mount_path, eml_path):
  try:
    logger.info("Executing eml_text function")
    path = config.get(mount_path, eml_path)
    files = glob.glob(path + '*.eml')
    logger.info("Number of email files from folder "+ path +" : "+str(len(files)))    
    for file in files:
      try:
        content = open(file, 'r').read()
        msg = email.message_from_string(content) 
        text = msg.get_payload()[0]
        name = file.split('/')    
        all_file_path = config.get(mount_path,'mnt_mpm2019_all_files')
        if not os.path.exists(all_file_path):
           path_exists(all_file_path)
        name = all_file_path +(name[-1].split('.'))[0]        
        eml_data = name + '.txt'       
        with open(eml_data, 'wb') as f:
          f.write(text.get_payload(decode=True))
      except Exception as e:
        logger.error("Error in eml_text iteration function :",exc_fino=True)
        logger.error("Error while extracting text from email:",file)
  except Exception as e:
    logger.error("Error in eml_text function",exc_info=True)
          
          
def eml_attachment(staging_path,eml_list,staging_path_pdf,raw_files,raw_format,all_files,file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor,source_type):
  try:
    logger.info("Executing eml_attachment function")   
    logger.info("Number of email files from folder : "+str(len(eml_list)))
    pdf_list_email = []
    for file in eml_list:
      try:
        content = open(file, 'r').read()
        msg = email.message_from_string(content)             
        attachment = msg.get_payload()[1]
        try:
          email_text = msg.get_payload()[0]
          email_name = file.split('/')[-1].rsplit('.')[0]
          eml_data = all_files +  email_name + '.txt'     
          if type(email_text) is not str:
            with open(eml_data, 'wb') as f:
              f.write(email_text.get_payload(decode=True))
            if eml_data.replace('//','/').strip() not in file_processing_blob_all_txt_list:
              file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, \
              {},'{}')".format(source_type, email_name, 
              'email','.eml', staging_path.replace('//','/'), eml_data.replace('//','/').strip(), 1,0,'GETDATE()','GETDATE()','null')
            else:
              file_processing_info_query = update_file_processing_info.format('GETDATE()',eml_data.replace('//','/').strip())
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.info('{}  email_text sucessfully'.format(eml_data))
        except Exception as e:
        #**************************************************************************************************************
        #Creation of insert query for the extracted invalid file path to the file_processing_info table and executed using
        #update_operation
        #***************************************************************************************************************
          file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} \
          ,'{}')".format(source_type, email_name, 
          'email','.eml', staging_path.replace('//','/'), 'null', 0,0,'null','null','null')
          update_operation(file_processing_info_query,sql_conn,cursor)
          logger.error('Error in email_text while processing {}'.format(eml_data))

        print('attachment',type(attachment))
        if type(attachment) is not str:
          filename = attachment.get_filename()
          if filename.endswith('.pdf'):#pdf_list_email
            file_copy_loc = staging_path_pdf
            f = open(file_copy_loc + filename, 'wb')
            f.write(attachment.get_payload(decode=True))
            pdf_list_email.append(file_copy_loc + filename)
            raw_files.append(file_copy_loc + filename)
            raw_format.append('.pdf')
            f.close()
            logger.info('{} sucessfully written in the path {}'.format(filename,file_copy_loc))
          else:
            logger.info('{} attachment from email not in PDF format so we are not processing'.format(filename))
        
      except Exception as e:
        logger.error("Error in eml_attachment iteration function :",exc_fino=True)
        logger.error("Error while extracting attachment from email:",file)
    logger.info('{} found in eml_list'.format(len(pdf_list_email)))
    return pdf_list_email 
  except Exception as e:
    logger.error("Error in outlook_attachment function",exc_info=True)

#************************************************************************************************************************************** 
#Function name: outlook_attachment
#Ojective : To fetch the attachments from outlook message 
#input parameters:
#msg_list : will hold outlook mesaage file folder path of respective sources
#staging_path: will hold the staging folder path of respective sources
#raw_files: its a list which will store the raw file path location 
#raw_format: its a list which will store the raw file format
#Usage: common code is written to fetch attachments from outlook message and write into the respective staging pdf raw path 
#called by : external_folder_structure_process
#**************************************************************************************************************************************     
def outlook_attachment(msg_list,staging_path_pdf,raw_files,raw_format):
  try:
    logger.info("Executing outlook_attachment function")   
    logger.info("Number of outlook files from folder : "+str(len(msg_list)))
    pdf_list_outlook = []
    for file in msg_list:
      try:
        with open(file) as msg_file:
            msg = Message(msg_file)
            attach = msg.attachments  
        for i in attach:
          with i.open() as attachment_fp:
              if i.filename.endswith('.pdf'):
                  file_copy_loc = staging_path_pdf
                  pdf_list_outlook.append(file_copy_loc + i.filename)
                  raw_files.append(file_copy_loc + i.filename)
                  raw_format.append('.pdf')
                  logger.info('{} file found in outlook_attachment'.format(i.filename))
                  with open(file_copy_loc + i.filename, 'wb') as my_data: ## write to temporary pdf file
                    my_data.write(attachment_fp.read())
                  logger.info('{} sucessfully written in the path {}'.format(i.filename,file_copy_loc))
              else:
                logger.info('{} attachment from outlook not in PDF format so we are not processing'.format(i.filename))
               
      except Exception as e:
        logger.error("Error in outlook_attachment iteration function :",exc_fino=True)
        logger.error("Error while extracting attachment from outlook:",file)
    logger.info('{} found in msg_list'.format(len(pdf_list_outlook)))
    return pdf_list_outlook 
  except Exception as e:
    logger.error("Error in outlook_attachment function",exc_info=True)
    
#*******************************************************************************************************************************    
#function name: extract_doc_text
#Objective: To convert documents to text files 
#input parameter:
#all_files: will hold all all-text folder path where extracted text files to be stored  
#staging_path: will hold the staging doc folder path of respective sources
#doc_file_list: will hold all the docx files in a list
#source_type: will hold the respective source type name
#file_processing_info: will hold the query to update extracted text file path  in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database 
#Usage: common code is written which extract text from a document files and store it in a text file on the respective sources
#called by : external_folder_structure_process
#*******************************************************************************************************************************
def extract_doc_text(staging_path,doc_file_list,source_type,all_files,file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor):
    for files in doc_file_list:
      try:
        doc = docx.Document(files)
        full_text = []
        for text in doc.paragraphs:
            full_text.append(text.text)
            text = '\n'.join(full_text)
        basenames=files.split('/')     
        file_name =  basenames[-1].rsplit('.',1)[0]
        basenames= all_files+basenames[-1].rsplit('.',1)[0]
        text_name = basenames.replace("/dbfs","dbfs:") + '.txt'            
        dbutils.fs.put(text_name,text,True)
        file_path = text_name.replace("dbfs:","/dbfs")
        #**************************************************************************************************************
        #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
        #update_operation
        #***************************************************************************************************************
        if file_path.replace('//','/').strip() not in file_processing_blob_all_txt_list:
          file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, \
          {},'{}')".format(source_type, file_name, 
        'Document','.docx', staging_path.replace('//','/'), file_path.replace('//','/').strip(), 1,0,'GETDATE()','GETDATE()','null')
        else:
          file_processing_info_query = update_file_processing_info.format('GETDATE()',file_path.replace('//','/').strip())
        update_operation(file_processing_info_query,sql_conn,cursor)
        logger.info('{}  extract_doc_text sucessfully'.format(files))
      except Exception as e:
        #**************************************************************************************************************
        #Creation of insert query for the extracted invalid file path to the file_processing_info table and executed using
        #update_operation
        #***************************************************************************************************************
        file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} \
        ,'{}')".format(source_type, file_name, 
        'Document','.docx', staging_path.replace('//','/'), 'null', 0,0,'null','null','null')
        update_operation(file_processing_info_query,sql_conn,cursor)
        logger.error('Error in extract_doc_text while processing {}'.format(files))
        
#********************************************************************************************
#This functionality extract text from a powerpoint files and store it in a text file.
#********************************************************************************************
def extract_pptx_text():
    try:
        ppt_nfiles = glob.glob(config.get('path', 'ppt_files') + '*.pptx')
        for files in ppt_nfiles:
            ppt = pptx.Presentation(files)
        
            text_runs = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
                            
            ppt_basenames = os.path.splitext(files)[0].replace('\\', '/').split('/')[-1]
            ppt_name = ppt_basenames + '.txt'                        
            f = open(config.get('path', 'extract_path') + ppt_name, 'a+', encoding='utf-8')
            f.write(text_runs)
            f.close() 
    except Exception as e:
        logger.error(e)

#*****************************************************************************************
#function name: sharepoint_native_scanned_pdf_split
#Ojective: Differentiating pdf into Native and Scanned
#staging_folder_path: will hold the staging pdf folder path of respective sources
#raw_pdf_files: will hold all the pdf files in a list
#Usage: Common code is written to split the raw pdf into native and scanned
#called by : external_folder_structure_process 
#******************************************************************************************
def sharepoint_native_scanned_pdf_split(staging_folder_path,raw_pdf_files):
    try:
      count = 0
      scan_files=[]
      native_files=[]
      logger.info("Executing sharepoint_native_scanned_pdf_split function")
      for files in raw_pdf_files:
          try:
              pdfFileObj = open(files, 'rb') 
              #********************************
              # creating a pdf reader object 
              #********************************
              pdfReader = PyPDF2.PdfFileReader(pdfFileObj)  
              if pdfReader.numPages>1:
                  pageObj = pdfReader.getPage(1)
              else:
                  pageObj = pdfReader.getPage(0) 

              if pageObj['/Resources'].get('/XObject') is not None:
                  scan_files.append(files)

              else:
                  native_files.append(files)
              pdfFileObj.close()
          except :
              scan_files.append(files)
              pdfFileObj.close()
      #**********************************************************************************************
      #scan_files: will hold all scanned pdf file path found in the staging folder 
      #copy_files: function will help to copy sacnned pdf from raw folder the scanned-files folder 
      #**********************************************************************************************
      if len(scan_files)>0:        
        logger.info('{} Number of scanned pdf files detected in: {}'.format(len(scan_files),staging_folder_path))
        staging_pdf_scanned = staging_folder_path.rsplit('/',2)[0] + '/scanned-files/'
        copy_files(scan_files, staging_pdf_scanned)
      else:
        staging_pdf_scanned = None
        logger.info('{} Number of scanned pdf files detected in: {}'.format(len(scan_files),staging_folder_path))
        
      #******************************************************************************
      #native_files: will hold all native pdf file path found in the staging folder
      #copy_files: function will help to copy native pdf from raw folder the native-files folder 
      #******************************************************************************  
      if len(native_files)>0:
        logger.info('{} Number of native pdf files detected in: {}'.format(len(native_files),staging_folder_path))
        staging_pdf_native = staging_folder_path.rsplit('/',2)[0] + '/native-files/'
        copy_files(native_files, staging_pdf_native)
      else:
        staging_pdf_native =None
        logger.info('{} Number of native pdf files detected in: {}'.format(len(native_files),staging_folder_path))      
      return staging_pdf_native, staging_pdf_scanned
    except Exception as e:
      logger.error("Error in sharepoint_native_scanned_pdf_split",exc_info=True)
      
#****************************************************************************************
#function name: intialize_temp_files
#Ojective: Temp folder for temporary execution
#Usage: TO create temporary folders for storing images while converting PDF to Image
#****************************************************************************************
def intialize_temp_files(temp_path):
    try:
        count=0
        logger.info("Executing intialize_temp_files function")
        temp = glob.glob(temp_path + '*.*')  
        if len(temp)==0:
            pass
        else:
            for i in temp:
              i = i.replace("/dbfs","dbfs:")
              dbutils.fs.rm(i)
              count+=1
        logger.info("Number of files got deleted from temporary folder : "+str(count))
    except Exception as e:
        logger.error("Error in initializing temp files function",exc_info=True)

#*****************************************************************************************************************************
#function name: pdf_to_image_converison
#Objectiv: To convert pdf to image
#input Parameter:
#files: will hold the pdf path which need to be converted into image 
#Ouput parameter:
#target: will hold the ouptut path where converted images will get stored 
#Usage: Common code is written to convert all the pages in the pdf to image in temporary location for tesseract processing
#called by: chemical_structure
#*****************************************************************************************************************************
def pdf_to_image_converison(files,target):
  try:
    logger.info("Executing pdf_to_image_converison function")
    destination=target
    if not os.path.exists(destination):
      os.mkdir(destination)
    with wimage(filename=files, resolution=300) as img:
       print('image')
       img.units = 'pixelsperinch'
       img.compression_quality = 99 
       img.save(filename = destination + '1.png')  
    logger.info("PDF file "+files+" has been converted into image file successfully")
  except Exception as e:
    logger.error("Error in pdf_to_image_converison",exc_info=True)
    logger.error("Error in image file",files)
    
#************************************************************************************************************************************** 
#Function name: native_pdf_extract_text
#Ojective : native pdf files  into text files    
#input parameters:
#native_path : will hold native pdf folder path of respective sources
#all_files: will hold all all-text folder path where extracted text files to be stored  
#staging_path: will hold the staging folder path of respective sources
#source_type: will hold the respective source type name
#file_processing_info: will hold the query to update extracted text file path  in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database  
#Usage: To extract text from scanned pdf files and stores text in a text file on the respective sourcs and insert the text file path
#into file_processing_info table
#called by : external_folder_structure_process
#**************************************************************************************************************************************   
def native_pdf_extract_text(native_path,all_files,staging_path,source_type,file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor):
    try:
        count=0
        logger.info("Executing native_pdf_extract_text function")
        native_files = glob.glob(native_path + '*.*')  
        logger.info("{} Number of native pdf files from folder {}".format(len(native_files),native_path))
        for files in native_files:
          try:
            text=''
            pdf_file = fitz.open(files)
            n_pages = pdf_file.pageCount
            for n in range(n_pages):
                page = pdf_file.loadPage(n)
                text = text + page.getText()
            basenames=files.split('/')     
            file_name =  basenames[-1].rsplit('.',1)[0]
            basenames= all_files+basenames[-1].rsplit('.',1)[0]
            text_name = basenames.replace("/dbfs","dbfs:") + '.txt'            
            dbutils.fs.put(text_name,text,True)
            file_path = text_name.replace("dbfs:","/dbfs")
            #**************************************************************************************************************
            #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
            #update_operation
            #***************************************************************************************************************
            if file_path.replace('//','/').strip() not in file_processing_blob_all_txt_list:
              file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} \
              ,'{}')".format(source_type, file_name,      
            'PDF','.pdf', staging_path.replace('//','/'), file_path.replace('//','/').strip(), 1,0,'GETDATE()','GETDATE()', 'null')
            else:
              file_processing_info_query = update_file_processing_info.format('GETDATE()',file_path.replace('//','/').strip())
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.info("Successfully extracted {} and updated the file_processing_info table".format(file_name))
            count+=1
          except Exception as e:
          #**************************************************************************************************************
          #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
          #update_operation
          #***************************************************************************************************************
            
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} \
            ,'{}')".format(source_type, file_name,     
            'PDF','.pdf', staging_path.replace('//','/'), 'Null', 0,0,'null','null','null')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.error("Error in native_pdf_extract_text function : iteraion",exc_info=True)
            logger.error("Error while extracting text from native file : {}".format(file_name))
        logger.info("Number of native pdf files got converted into text files successfully : "+str(count))
    except Exception as e:
        logger.error("Error in native_pdf_extract_text function",exc_info=True)
        
#************************************************************************************************************************************** 
#Function name: scanned_pdf_extract_text
#Ojective : Scanned pdf files  into text files  
#input parameters:
#scanned_path : will hold scanned pdf folder path of respective sources
#all_files: will hold all all-text folder path where extracted text files to be stored  
#staging_path: will hold the staging folder path of respective sources
#source_type: will hold the respective source type name
#file_processing_info: will hold the query to update extracted text file path  in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#Usage: To extract text from scanned pdf files and stores text in a text file on the respective sourcs and insert the text file path
#into file_processing_info table
#called by : external_folder_structure_process
#**************************************************************************************************************************************       
def scanned_pdf_extract_text(scanned_path,all_files,staging_path,source_type,file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor):
    try: 
        logger.info("Executing scanned_pdf_extract_text function")
        scanned_files = glob.glob(scanned_path + '*.pdf')
        logger.info("Number of scanned pdf files from folder "+scanned_path+" : "+str(len(scanned_files)))
        count=0
        temp = all_files.rsplit('/',2)[0] + '/temp/'
        for files in scanned_files:
          try:
            #intialize_temp_files(path, temp)
            #intialize_temp_files was replace by temp function
            path_exists(temp)
            #calling pdf to image conversion function
            pdf_to_image_converison(files,temp)
            image_files = glob.glob(temp + '*')
            text_extract = ''
            for j in range(len(image_files)):
                im = PIL.Image.open(image_files[j])
                if im.mode=='P':
                    im = im.convert(palette=0)
                im1 = im.filter(ImageFilter.EDGE_ENHANCE_MORE)                                    
                config1 = (' --psm 6')
                text_val = pyt.image_to_string(im1, config=config1)          
                text_extract = text_extract + text_val
            basenames=files.split('/')       
            file_name =  basenames[-1].rsplit('.',1)[0]
            basenames=all_files+(basenames[-1].rsplit('.',1))[0]
            text_name = basenames.replace("/dbfs","dbfs:") + '.txt'            
            dbutils.fs.put(text_name,text_extract,True)
            file_path = text_name.replace("dbfs:","/dbfs")
            #************************************************************************************************************************
            #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
            #update_operation
            #************************************************************************************************************************
            if file_path.replace('//','/').strip() not in file_processing_blob_all_txt_list:
              file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}' ,{}, \
              {},'{}')".format(source_type, file_name,     
            'PDF','.pdf',staging_path.replace('//','/'), file_path.replace('//','/').strip(), 1,0,'GETDATE()','GETDATE()','null')
            else:
              file_processing_info_query = update_file_processing_info.format('GETDATE()',file_path.replace('//','/').strip()) 
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.info("Successfully extracted {} and updated the file_processing_info table".format(file_name))
            count+=1
          except Exception as e:
          #************************************************************************************************************************
          #Creation of insert query for the extracted invalid file path to the file_processing_info table and executed using
          #update_operation
          #************************************************************************************************************************
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} \
            ,'{}')".format(source_type, file_name,     
            'PDF','.pdf', staging_path.replace('//','/'), 'Null', 0,0,'null','null','null')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.error("Error in scanned_pdf_extract_text function : iteration",exc_info=True)
            logger.error("Error while extracting text from scanned file : {}".format(file_name))
        logger.info("Number of scanned pdf files got converted into text files successfully : "+str(count))
    except Exception as e:
        logger.error("Error in scanned_pdf_extract_text function",exc_info=True)
        
        
#**************************************************************************************************************************************
#Function name: external_folder_structure_process
#objective: To convert all the file types into text format 
#input Parameters:
#external_folder_structure: Will call external_source_data function by passing external_folder_structure_query and returns all the
#                           details in the external_folder_structure in a dataframe
#external_source_file_formats: Will call external_source_data function by passing file_format_query and returns all the
#                              inscope file formats to be consider for processing
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database                
#Usage: Extraction of input file data fetched from the external_folder_structure table and stores all the extracted file path in the
#file_processing_info table
#ouput: returns raw_df which holds all the staging file path in dataframe which helps to move file to processed folder
#**************************************************************************************************************************************
def external_folder_structure_process(external_folder_structure,external_source_file_formats,file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor):
  try:
    #raw_df will be used for moving the raw files into processed folder after the key-data extract
    #raw files and raw format will conatin each raw file path and formats in the list
    #after all the file extraction raw files and raw format will append to the raw_df dataframe
    raw_df1 = pd.DataFrame()
    raw_files = []
    raw_format = []  
    for index in external_folder_structure.index:
      source_type = external_folder_structure['blob_src_type'][index].strip()
      mount_path = external_folder_structure['db_fs_mount_path'][index].strip()
      staging_path =  (mount_path + external_folder_structure['absolute_path'][index]).replace('//','/')
      print(staging_path)
      try:
        logger.info('Text extraction started for {}'.format(staging_path))
        if os.path.exists(staging_path):
           # print('exist')
            all_files = staging_path.split('staging',1)[0] + 'analytics/processed/all-text/'
            excel_files = staging_path.split('staging',1)[0] + 'analytics/processed/excel/'
            if not os.path.exists(all_files):
              path_exists(all_files)
            if not os.path.exists(excel_files):
              #print()
              path_exists(excel_files)
            #***************************************************************************
            #fetching all the pdf file types from the sources
            #pdf_file_list: will have all the pdf file path from each category
            #***************************************************************************
            print(external_source_file_formats)
            if '.pdf' in external_source_file_formats:
                pdf_file_list = glob.glob(staging_path +'*.pdf')     
            #    pdf_file_list = ['/dbfs/mnt/momentive-sources-pih/sharepoint-pih/customer-communications-pih/mpm-2019-pih/staging/pdf/raw/Polydimethylsiloxane - Review 6974.pdf']
                                  
                if bool(pdf_file_list):
                    logger.info('{} pdf files found in the {}'.format(len(pdf_file_list),staging_path))
                    raw_files = raw_files + pdf_file_list
                    raw_format = raw_format + ['.pdf']*len(pdf_file_list)
                    #***************************************************************************************************
                    #sharepoint_native_scanned_pdf_split: will split the pdf files into two types like(native, scanned)
                    #***************************************************************************************************
#                     pdf_file_list1=[]
#                     for i in pdf_file_list:
#                       pdf_check1='Silopren LSR 7180 Marketing Bulletin'
#                       if pdf_check1 in i:
#                         pdf_file_list1.append(i) 
#                     raw_files = raw_files + pdf_file_list1
#                    raw_format = raw_format + ['.pdf']*len(pdf_file_list1)
                    native_path, scanned_path = sharepoint_native_scanned_pdf_split(staging_path,pdf_file_list)
                    if native_path != None:
                     #***************************************************************************************************
                     #native_pdf_extract_text: will extract data from the native pdf type
                     #**************************************************************************************************
                        native_pdf_extract_text(native_path,all_files,staging_path,source_type,file_processing_info,
                        update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)
                    if scanned_path != None:
                     #***************************************************************************************************
                     #scanned_pdf_extract_text: will extract data from the scanned pdf type
                     #***************************************************************************************************
                        scanned_pdf_extract_text(scanned_path,all_files,staging_path,source_type,file_processing_info,
                        update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)
          #**************************************************************
          #fetching all the Document file types from the sources
          #**************************************************************        
            if '.docx' in  external_source_file_formats:
                doc_file_list = glob.glob(staging_path+'*.docx')
                if bool(doc_file_list):
                    logger.info('{} docx files found in the {}'.format(len(doc_file_list),staging_path))
                    raw_files = raw_files + doc_file_list
                    raw_format = raw_format + ['.docx']*len(doc_file_list)
                    #***************************************************************************************************
                    #extract_doc_text: will extract data from the documnet file type
                    #doc_file_list: will have all the document file path from each category
                    #***************************************************************************************************
                    extract_doc_text(staging_path,doc_file_list,source_type,all_files,file_processing_info,
                    update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)
          #******************************************************************************
          #fetching all the message file types from the sources
          #msg_list: will have all the message file path from each category
          #*******************************************************************************      
#             if '.msg' in  external_source_file_formats:
#                 msg_list = glob.glob(staging_path+'*.msg')
#                 staging_path_pdf = staging_path.lower()+'staging/pdf/raw/' 
#                 path_exists(staging_path_pdf)
#                 if bool(msg_list):  
#                    #*********************************************************************
#                    #outlook_attachment: will fetch the attachments found in the messasge
#                    #*********************************************************************              
#                    pdf_out_look = outlook_attachment(msg_list,staging_path_pdf,raw_files,raw_format)
#                    raw_files = raw_files + pdf_out_look
#                    raw_format = raw_format + ['.pdf']*len(pdf_out_look)
#                    if bool(pdf_out_look):
#                      native_path, scanned_path = sharepoint_native_scanned_pdf_split(staging_path_pdf,pdf_out_look)                       
#                      if native_path != None:
#                        native_pdf_extract_text(native_path,all_files,staging_path_pdf,source_type,file_processing_info,
#                        update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)
#                      if scanned_path != None:
#                        scanned_pdf_extract_text(scanned_path,all_files,staging_path_pdf,source_type,file_processing_info,
#                        update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor) 

          #******************************************************************************
          #fetching all the message file types from the sources
          #eml_list: will have all the message file path from each category
          #*******************************************************************************      
#             if '.eml' in  external_source_file_formats:
#                 eml_list = glob.glob(staging_path+'*.eml')
#                 staging_path_index = staging_path.lower().find('staging')
#                 if staging_path_index != -1:
#                   staging_path_pdf = staging_path[:staging_path_index] + 'staging/pdf/raw/'
#                 if bool(eml_list):  
#                    #*********************************************************************
#                    #outlook_attachment: will fetch the attachments found in the messasge
#                    #*********************************************************************              
#                    pdf_email= eml_attachment(staging_path,eml_list,staging_path_pdf,raw_files,raw_format,all_files,file_processing_info,
#                               update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor,source_type)
#                    raw_files = raw_files + pdf_email
#                    raw_format = raw_format + ['.pdf']*len(pdf_email)
#                    if bool(pdf_email):
#                      native_path, scanned_path = sharepoint_native_scanned_pdf_split(staging_path_pdf,pdf_email) 
#                      if native_path != None:
#                        native_pdf_extract_text(native_path,all_files,staging_path_pdf,source_type,file_processing_info,
#                        update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)
#                      if scanned_path != None:
#                        scanned_pdf_extract_text(scanned_path,all_files,staging_path_pdf,source_type,file_processing_info,
#                        update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor) 
          #**************************************************************
          #fetching all the csv file types from the sources
          #csv_list: will have all the csv format file path from each category
          #************************************************************** 
            if '.csv' in external_source_file_formats:
                csv_list = glob.glob(staging_path+'*.csv')
                raw_files = raw_files + csv_list
                raw_format = raw_format + ['.csv']*len(csv_list)
                #print(csv_list)
                if bool(csv_list):
                  #*********************************************************************
                  #csv_text_extract: will extract the data from the csv file type
                  #*********************************************************************
                  logger.info('{} csv file found in the staging_path'.format(len(csv_list)))
                  csv_text_extract(staging_path,csv_list,source_type,all_files,excel_files,file_processing_info,
                  update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)
                                                     
          #*******************************************************************************
          #fetching all the xlsx and xlsm file types from the sources
          #xlsx_list: will have all the xlsx format file path from each category
          #xlsm_list: will have all the xlsm format file path from each category
          #*******************************************************************************
            if 'xlsx'  or 'xlsm' in external_source_file_formats:
                xlsx_list = glob.glob(staging_path+'*.xlsx')       
                xlsm_list = glob.glob(staging_path+'*.xlsm')
                raw_files = raw_files + xlsx_list
                raw_format = raw_format + ['.xlsx']*len(xlsx_list)
                raw_files = raw_files + xlsm_list
                raw_format = raw_format + ['.xlsm']*len(xlsm_list)
                xlsx_list = xlsx_list + xlsm_list
                #print(xlsx_list)
#                 xlsx_list1=[]
#                 for i in xlsx_list:
#                   pdf_check1='Selant - Monthly update 1-9-2018'
#                   if pdf_check1 in i:
#                     xlsx_list1.append(i)
                #****************************************************************************
                #xlsx_text_extract: will extract the data from the xlsx and xlsm file type
                #****************************************************************************
                #print('xlsx_list1',xlsx_list1)
                if bool(xlsx_list):
                    logger.info('{} excel file found in the staging_path'.format(len(xlsx_list)))
                    xlsx_text_extract(staging_path,xlsx_list,source_type,all_files,excel_files,file_processing_info,
                    update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)
        else:
          logger.info('Invalid staging path {}'.format(staging_path))
      except Exception as e:
        logger.error('Something went wrong while text extraction of {} folder'.format(staging_path), exc_info=True)             
    raw_df1['file_name'] = raw_files 
    return raw_df1
  except Exception as e:
    logger.error('Something went wrong in the external_folder_structure_process function', exc_info=True)
    
#*************************************************************************************************************************************
#function name: update_operation
#Objective: insert, update and Delete operations on the table
#Usage: common function is written to perform (insert, update and Delete) query operation on the required table 
#input Parameters: sql_conn will hold DB_connectivity object and Cursor will hold cursor object for executing queries, it helps to 
#                  update the database
#called by: native_pdf_extract_text, scanned_pdf_extract_text, extract_doc_text, xlsx_text_extract, csv_text_extract,file_validation #functions   
#*************************************************************************************************************************************

def update_operation(query,sql_conn,cursor):
  try:
    all_text_find = query[query[:query.find('.txt')].rfind('/dbfs'):query.find('.txt')+4]
    print(all_text_find)
    if all_text_find:    
      extracted_file_list.append(all_text_find)
    cursor.execute(query)
    sql_conn.commit()  
    logger.info('Successfully executed the query {}'.format(query), exc_info=True)
  except Exception as e:
    logger.error('Something went wrong in the update_operation while executing the query {}'.format(query), exc_info=True)
#********************************************************************************************************************************
#function name: external_source_data
#Objective:Select operations on the table
#input Parameters: sql_conn will hold DB_connectivity object and query will hold select operations or query to fetch the data from 
#                  the table
#ouput Parameter: result will hold the fetched data from azure sql table in a dataframe
#called by: main, pattern_match_validation function
#Usage: common function is written to perform select query operation on the required table and returns the result in dataframe
#********************************************************************************************************************************     

def external_source_data(sql_conn,query):
  try:
    if sql_conn is not None:  
      result = pd.read_sql(query, sql_conn) 
      logger.info('Successfully executed the query {}'.format(query), exc_info=True) 
    else:
      logger.error('Sql_conn has None value something went wrong in the Sql server connection')     
    return result
  except Exception as e:
    logger.error('Something went wrong in the external_source_data function while executing the query {}'.format(query), exc_info=True)
    
#**********************************************************************************************
#function name: Sql_db_connection
#Objective: connecting sql db using pyodbc
#Usage: common function is written to connect with given database using pyodbc package
#output: Sql_conn will hold the DB_connectivity object
#called by : Main function
#**********************************************************************************************

def Sql_db_connection(): 
  try:
    logger.info('Connecting azure sql server')
    server = config.get('sql_db', 'server')
    database = config.get('sql_db', 'database')
    username = config.get('sql_db', 'username')
    password = config.get('sql_db', 'password')
    DATABASE_CONFIG = {'server': server,'database': 'cld-it-dev-pih-db1','username': username,'password': password}
    driver= "{ODBC Driver 17 for SQL Server}"
    connection_string = 'DRIVER=' + driver + \
                      ';SERVER=' + DATABASE_CONFIG['server'] + \
                      ';PORT=1433' + \
                      ';DATABASE=' + DATABASE_CONFIG['database'] + \
                      ';UID=' + DATABASE_CONFIG['username'] + \
                      ';PWD=' + DATABASE_CONFIG['password'] 
    sql_conn = pyodbc.connect(connection_string)
    logger.info('Successfully connected with the Azure sql serevr ')
    if sql_conn is None:
      logger.error('sql is not connected properly is returns None Object ')    
    return sql_conn    
  except Exception as e:
    logger.error('Something went wrong in the Sql_db_connection function', exc_info=True)

#****************************************************************************************************************************************** 
#Function name: Main 
#Objective: Program will start process using this function 
#sql_conn: Azure SQl DB Connectivity will be created using this  Sql_db_connection()   
#cursor: Cursor will be created using this sql_conn.cursor() for executing Sql operations
#external_source_folder_structure table: external_source_folder_structure table will contain all the data ingestion details from azure to 
#                                  blob storgae based on each category  
#external_folder_structure_query: Will have the query to connect with  external_source_folder_structure table
#file_format table: will have all the inscope file formats(like 'PDF', "Document') etc to be processed
#file_format_query: inscope file format query will get capture in this variable
#file_processing_info table: will have all the extracted file path details, Sources type (like Sharepoint,Website etc) and  
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#external_source_data: Will perform sql select operation by passing two parameter DB Connectivity(sql_conn) and select query.
#                      it will return output in dataframe 
#external_folder_structure: Will call external_source_data function by passing external_folder_structure_query and returns all the
#                           details in the external_folder_structure in a dataframe
#external_source_file_formats: Will call external_source_data function by passing file_format_query and returns all the
#                              inscope file formats to be consider for processing
#external_folder_structure_process: its a function which takes 5 parameters like(external_folder_structure,external_source_file_formats,
#                                   file_processing_info, ,sql_conn, sql_conn) using this will etract all the file data got from the  
#external_file_process_query: 
#pattern_match_validation: it will do the pattern matching for 15 categories for extracted files using file information stored in the    
#                           file_processing_info table and pattern matching keywords stored in the pattern_matching_keys table
#*****************************************************************************************************************************************  
   
def main(): 
  try:
      logger.info('Executing main function')
      sql_conn = Sql_db_connection()
      cursor = sql_conn.cursor()  
      external_folder_structure_query = config.get('mount_path', 'external_source_folder_structure')
      file_format_query = config.get('mount_path', 'external_source_file_formats')
      file_processing_blob_all_txt_info = config.get('mount_path', 'file_processing_blob_all_txt_info')
      file_processing_info = config.get('mount_path', 'file_processing_info')
      update_file_processing_info = config.get('mount_path', 'update_file_processing_info')
      external_folder_structure = external_source_data(sql_conn,external_folder_structure_query)
      external_source_file_formats = external_source_data(sql_conn,file_format_query)['file_format'].values.tolist()
      file_processing_blob_all_txt_list = external_source_data(sql_conn,file_processing_blob_all_txt_info)\
                                          ['blob_all_txt_file_path'].values.tolist()
      raw_df = external_folder_structure_process(external_folder_structure,external_source_file_formats,
               file_processing_info,update_file_processing_info,file_processing_blob_all_txt_list,sql_conn,cursor)           
      external_file_process_query = config.get('mount_path', 'external_file_process')
      external_processed_files_df = external_source_data(sql_conn,external_file_process_query)
      unstruct_category_key_query = config.get('mount_path','unstruct_category_key')
      unstruct_category_key_df = external_source_data(sql_conn,unstruct_category_key_query)
      unstructure_processed_data_query = config.get('mount_path', 'unstructure_processed_data')
      pattern_match_validation(sql_conn,external_processed_files_df,cursor,unstruct_category_key_df
                              ,raw_df,unstructure_processed_data_query)
      valid_folder_list =[]                       
  except Exception as e:
    logger.error('Something went wrong in main function', exc_info=True)
    
#***************************************************************************************** 
#calling the main function when python code is triggered from azure pipeline                                      
#***************************************************************************************** 

if __name__ == '__main__':
  try:
    logger.info('calling main function')
    main()
  except Exception as e:
    logger.error('Somethng went wrong while calling main function',exc_info=True)

# COMMAND ----------

# MAGIC %sh
# MAGIC cat /databricks/driver/shared_main_code_error2.log

# COMMAND ----------

consol_data.head()

# COMMAND ----------

path_exists('/dbfs/mnt/momentive-sources-pih/sharepoint-pih/tox-study-pih/sealants-silanes-library-pih/analytics/processed/excel/Silane Monthly Tox update 9-26-2018/'+'relevant_data_files/')
final.to_csv('/dbfs/mnt/momentive-sources-pih/sharepoint-pih/tox-study-pih/sealants-silanes-library-pih/analytics/processed/excel/Silane Monthly Tox update 9-26-2018/'+'relevant_data_files/' + 'non_relevant_data' +'.csv',index=None, header=True)


# COMMAND ----------

import shutil
shutil.rmtree('/dbfs/mnt/momentive-sources-pih/sharepoint-pih/toxicology-pih/tir-sip-pih/analytics/'.replace('//','/'))

# COMMAND ----------

# import shutil
# shutil.rmtree('/dbfs/mnt/momentive-sources-pih/sharepoint-pih/tox-study-pih/sealants-silanes-library-pih/analytics/'.replace('//','/'))

# COMMAND ----------

a='nxt'
rgx_img = re.compile('{}'.format('NXT'.replace('*','\*')),re.I)
for match in re.finditer(rgx_img,a):
  match_f = match.group()
  print(match_f)

# COMMAND ----------

import datetime
print(str(datetime.datetime.now())[:19])

# COMMAND ----------

import shutil
shutil.rmtree('/dbfs/mnt/momentive-sources-pih/website-pih/product-application-pih/analytics/'.replace('//','/'))

# COMMAND ----------
